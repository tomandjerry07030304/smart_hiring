"""
🏗️ SYSTEM ARCHITECTURE SLIDE - Clean Professional Design
Single slide for easy copy-paste
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK_BG = RGBColor(45, 45, 48)
TITLE_GOLD = RGBColor(212, 175, 55)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(200, 200, 200)

# Layer colors - gradient blues
LAYER_1 = RGBColor(41, 128, 185)   # Presentation Layer - Blue
LAYER_2 = RGBColor(39, 174, 96)    # Application Layer - Green
LAYER_3 = RGBColor(142, 68, 173)   # Service Layer - Purple
LAYER_4 = RGBColor(230, 126, 34)   # Data Layer - Orange

slide = prs.slides.add_slide(prs.slide_layouts[6])

# Background
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK_BG
bg.line.fill.background()

# Title
title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(12.733), Inches(0.7))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "System Architecture"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Underline
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(0.85), Inches(5.333), Inches(0.025))
line.fill.solid()
line.fill.fore_color.rgb = TITLE_GOLD
line.line.fill.background()

# ============================================
# LAYER 1: PRESENTATION LAYER (Top)
# ============================================
layer1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.1), Inches(12.333), Inches(1.15))
layer1.fill.solid()
layer1.fill.fore_color.rgb = LAYER_1
layer1.line.color.rgb = WHITE
layer1.line.width = Pt(1.5)

# Layer 1 Title
l1_title = slide.shapes.add_textbox(Inches(0.7), Inches(1.15), Inches(3), Inches(0.4))
tf = l1_title.text_frame
p = tf.paragraphs[0]
p.text = "PRESENTATION LAYER"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = WHITE

# Layer 1 Components
components_1 = ["Candidate Portal", "Recruiter Dashboard", "Admin Panel", "Fairness Audit UI"]
x_pos = 0.8
for comp in components_1:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(1.55), Inches(2.8), Inches(0.55))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(52, 152, 219)
    box.line.color.rgb = WHITE
    box.line.width = Pt(1)
    
    txt = slide.shapes.add_textbox(Inches(x_pos), Inches(1.62), Inches(2.8), Inches(0.45))
    tf = txt.text_frame
    p = tf.paragraphs[0]
    p.text = comp
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    x_pos += 3.05

# Arrow down
arrow1 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.2), Inches(2.3), Inches(0.9), Inches(0.35))
arrow1.fill.solid()
arrow1.fill.fore_color.rgb = TITLE_GOLD
arrow1.line.fill.background()

# ============================================
# LAYER 2: APPLICATION LAYER
# ============================================
layer2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.7), Inches(12.333), Inches(1.15))
layer2.fill.solid()
layer2.fill.fore_color.rgb = LAYER_2
layer2.line.color.rgb = WHITE
layer2.line.width = Pt(1.5)

l2_title = slide.shapes.add_textbox(Inches(0.7), Inches(2.75), Inches(3), Inches(0.4))
tf = l2_title.text_frame
p = tf.paragraphs[0]
p.text = "APPLICATION LAYER"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = WHITE

# Layer 2 Components
components_2 = ["Flask REST API", "JWT Authentication", "Rate Limiting", "CORS Handler"]
x_pos = 0.8
for comp in components_2:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(3.15), Inches(2.8), Inches(0.55))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(46, 204, 113)
    box.line.color.rgb = WHITE
    box.line.width = Pt(1)
    
    txt = slide.shapes.add_textbox(Inches(x_pos), Inches(3.22), Inches(2.8), Inches(0.45))
    tf = txt.text_frame
    p = tf.paragraphs[0]
    p.text = comp
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    x_pos += 3.05

# Arrow down
arrow2 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.2), Inches(3.9), Inches(0.9), Inches(0.35))
arrow2.fill.solid()
arrow2.fill.fore_color.rgb = TITLE_GOLD
arrow2.line.fill.background()

# ============================================
# LAYER 3: SERVICE LAYER
# ============================================
layer3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.3), Inches(12.333), Inches(1.15))
layer3.fill.solid()
layer3.fill.fore_color.rgb = LAYER_3
layer3.line.color.rgb = WHITE
layer3.line.width = Pt(1.5)

l3_title = slide.shapes.add_textbox(Inches(0.7), Inches(4.35), Inches(3), Inches(0.4))
tf = l3_title.text_frame
p = tf.paragraphs[0]
p.text = "SERVICE LAYER"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = WHITE

# Layer 3 Components - 5 services
components_3 = ["NLP Engine", "Fairness Engine", "Matching Engine", "Email Service", "Analytics"]
x_pos = 0.6
for comp in components_3:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(4.75), Inches(2.3), Inches(0.55))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(155, 89, 182)
    box.line.color.rgb = WHITE
    box.line.width = Pt(1)
    
    txt = slide.shapes.add_textbox(Inches(x_pos), Inches(4.82), Inches(2.3), Inches(0.45))
    tf = txt.text_frame
    p = tf.paragraphs[0]
    p.text = comp
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    x_pos += 2.45

# Arrow down
arrow3 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.2), Inches(5.5), Inches(0.9), Inches(0.35))
arrow3.fill.solid()
arrow3.fill.fore_color.rgb = TITLE_GOLD
arrow3.line.fill.background()

# ============================================
# LAYER 4: DATA LAYER
# ============================================
layer4 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.9), Inches(12.333), Inches(1.15))
layer4.fill.solid()
layer4.fill.fore_color.rgb = LAYER_4
layer4.line.color.rgb = WHITE
layer4.line.width = Pt(1.5)

l4_title = slide.shapes.add_textbox(Inches(0.7), Inches(5.95), Inches(3), Inches(0.4))
tf = l4_title.text_frame
p = tf.paragraphs[0]
p.text = "DATA LAYER"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = WHITE

# Layer 4 Components - with icons
components_4 = [
    ("🗄️ MongoDB", "Documents & Resumes"),
    ("⚡ Redis", "Cache & Sessions"),
    ("📁 File Storage", "Resume Files"),
    ("📋 Audit Logs", "Compliance Trail")
]
x_pos = 0.8
for icon_name, desc in components_4:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(6.3), Inches(2.8), Inches(0.6))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(243, 156, 18)
    box.line.color.rgb = WHITE
    box.line.width = Pt(1)
    
    txt = slide.shapes.add_textbox(Inches(x_pos), Inches(6.32), Inches(2.8), Inches(0.3))
    tf = txt.text_frame
    p = tf.paragraphs[0]
    p.text = icon_name
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    txt2 = slide.shapes.add_textbox(Inches(x_pos), Inches(6.55), Inches(2.8), Inches(0.3))
    tf = txt2.text_frame
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(255, 255, 220)
    p.alignment = PP_ALIGN.CENTER
    
    x_pos += 3.05

# ============================================
# Right side: Technology badges
# ============================================
tech_box = slide.shapes.add_textbox(Inches(11.5), Inches(1.6), Inches(1.5), Inches(0.3))
tf = tech_box.text_frame
p = tf.paragraphs[0]
p.text = "HTML/CSS/JS"
p.font.size = Pt(9)
p.font.color.rgb = LIGHT_GRAY
p.alignment = PP_ALIGN.RIGHT

tech_box2 = slide.shapes.add_textbox(Inches(11.5), Inches(3.2), Inches(1.5), Inches(0.3))
tf = tech_box2.text_frame
p = tf.paragraphs[0]
p.text = "Python Flask"
p.font.size = Pt(9)
p.font.color.rgb = LIGHT_GRAY
p.alignment = PP_ALIGN.RIGHT

tech_box3 = slide.shapes.add_textbox(Inches(11.5), Inches(4.8), Inches(1.5), Inches(0.3))
tf = tech_box3.text_frame
p = tf.paragraphs[0]
p.text = "spaCy/BERT"
p.font.size = Pt(9)
p.font.color.rgb = LIGHT_GRAY
p.alignment = PP_ALIGN.RIGHT

tech_box4 = slide.shapes.add_textbox(Inches(11.5), Inches(6.4), Inches(1.5), Inches(0.3))
tf = tech_box4.text_frame
p = tf.paragraphs[0]
p.text = "NoSQL/Cache"
p.font.size = Pt(9)
p.font.color.rgb = LIGHT_GRAY
p.alignment = PP_ALIGN.RIGHT

# Save
output_path = os.path.join(os.path.dirname(__file__), "Architecture_Slide_Clean.pptx")
prs.save(output_path)

print("=" * 50)
print("✅ CLEAN ARCHITECTURE SLIDE CREATED!")
print(f"📁 File: {output_path}")
print("=" * 50)
print("\n📋 Instructions:")
print("   1. Open Architecture_Slide_Clean.pptx")
print("   2. Select All (Ctrl+A) on the slide")
print("   3. Copy (Ctrl+C)")
print("   4. Go to your main presentation")
print("   5. Paste (Ctrl+V) on the target slide")
