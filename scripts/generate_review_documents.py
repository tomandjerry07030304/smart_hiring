#!/usr/bin/env python3
"""
Smart Hiring — Review Document Generator
==========================================
Generates comprehensive DOCX, PPTX, and PDF for academic/project review.
Covers: Architecture, AI/ML, Phases, Methodology, Interview Design, Roadmap.

Run: python generate_review_documents.py
Output: review_output/ folder
"""

import os
import sys
import math
import base64
import urllib.request
import urllib.error
from pathlib import Path
from datetime import datetime, timedelta
from io import BytesIO

# ── Third-party imports ──────────────────────────────
from docx import Document
from docx.shared import Inches, Pt, Cm, RGBColor, Emu
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn, nsdecls
from docx.oxml import parse_xml

from pptx import Presentation
from pptx.util import Inches as PI, Pt as PP, Emu as PE
from pptx.dml.color import RGBColor as PC
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from PIL import Image, ImageDraw, ImageFont

try:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import inch, cm
    from reportlab.lib.colors import HexColor, Color
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_JUSTIFY
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
        PageBreak, Image as RLImage, KeepTogether
    )
    HAS_REPORTLAB = True
except ImportError:
    HAS_REPORTLAB = False
    print("  [WARN] reportlab not installed — PDF will be skipped. Run: pip install reportlab")

# ═══════════════════════════════════════════════════════════════
#  CONFIGURATION
# ═══════════════════════════════════════════════════════════════
BASE_DIR = Path(r"c:\Users\venkat anand\OneDrive\Desktop\@smart\smart_hiring")
OUTPUT_DIR = BASE_DIR / "review_output"
DIAG_DIR = OUTPUT_DIR / "diagrams"

# Color palette
NAVY = (11, 31, 59)
DARK_BLUE = (0, 51, 102)
ROYAL_BLUE = (37, 99, 235)
STEEL = (70, 130, 180)
EMERALD = (16, 185, 129)
AMBER = (245, 158, 11)
RED_ACCENT = (239, 68, 68)
WHITE = (255, 255, 255)
LIGHT_GRAY = (249, 250, 251)
DARK_GRAY = (30, 41, 59)
MED_GRAY = (100, 116, 139)

PROJECT_TITLE = "Smart Hiring"
PROJECT_SUBTITLE = "AI-Powered Intelligent Recruitment Platform"
VERSION = "v4.2+"
TEAM_SIZE = "4-6 members"
DATE_STR = datetime.now().strftime("%B %d, %Y")


# ═══════════════════════════════════════════════════════════════
#  MERMAID DIAGRAM RENDERING
# ═══════════════════════════════════════════════════════════════
MERMAID_DIAGRAMS = {
    "system_architecture": """graph TD
    subgraph Frontend["Frontend (SPA)"]
        UI[Web Interface]
        Auth[Auth Module]
        Dash[Dashboard]
    end

    subgraph Backend["Backend (Flask API)"]
        API[REST API Layer]
        Routes[Blueprint Routes]
        Services[Service Layer]
        MW[Middleware]
    end

    subgraph Intelligence["AI/ML Engine"]
        SBERT[SBERT Embeddings]
        NLP[spaCy NLP Parser]
        RANK[Ranking Engine]
        FAIR[Fairness Engine]
    end

    subgraph Data["Data Layer"]
        MONGO[(MongoDB)]
        REDIS[(Redis Cache)]
    end

    subgraph Workers["Background Workers"]
        CELERY[Celery Workers]
        QUEUE[Task Queue]
    end

    UI --> API
    API --> Routes
    Routes --> Services
    Routes --> MW
    Services --> Intelligence
    Services --> MONGO
    Services --> REDIS
    API --> QUEUE
    QUEUE --> CELERY
    CELERY --> Services""",

    "hiring_pipeline": """graph LR
    JP[Job Posting] --> CS[Candidate Sourcing]
    CS --> RI[Resume Ingestion]
    RI --> NLP[NLP Parsing]
    NLP --> ML[AI Screening]
    ML --> ASS[Assessment]
    ASS --> INT[Interview]
    INT --> DE[Decision Engine]
    DE --> OFF[Offer Management]
    OFF --> ONB[Onboarding]""",

    "ai_pipeline": """graph TD
    subgraph Ingestion
        PDF[PDF/DOCX Upload]
        EXTRACT[Text Extraction]
        CLEAN[Text Cleaning]
    end

    subgraph NLP_Pipeline["NLP Processing"]
        SPACY[spaCy NER]
        SKILLS[Skill Extraction]
        EXP[Experience Parsing]
        EDU[Education Detection]
    end

    subgraph Matching["Semantic Matching"]
        SBERT[SBERT all-MiniLM-L6-v2]
        TFIDF[TF-IDF Fallback]
        COSINE[Cosine Similarity]
    end

    subgraph Scoring["Scoring Engine"]
        WEIGHT[Weighted Aggregation]
        CCI[Career Consistency Index]
        EXPLAIN[Explainability Layer]
    end

    subgraph Fairness["Fairness Layer"]
        BIAS[Bias Detection]
        ANON[Anonymization]
        AUDIT[Audit Trail]
    end

    PDF --> EXTRACT --> CLEAN
    CLEAN --> SPACY
    SPACY --> SKILLS
    SPACY --> EXP
    SPACY --> EDU
    SKILLS --> SBERT
    SBERT --> COSINE
    TFIDF --> COSINE
    COSINE --> WEIGHT
    EXP --> WEIGHT
    EDU --> WEIGHT
    WEIGHT --> CCI
    CCI --> EXPLAIN
    EXPLAIN --> BIAS
    BIAS --> ANON
    ANON --> AUDIT""",

    "state_machine": """stateDiagram-v2
    [*] --> APPLIED
    APPLIED --> SCREENING
    SCREENING --> SHORTLISTED
    SCREENING --> REJECTED
    SHORTLISTED --> ASSESSMENT_PENDING
    ASSESSMENT_PENDING --> ASSESSMENT_COMPLETED
    ASSESSMENT_COMPLETED --> INTERVIEW_SCHEDULED
    INTERVIEW_SCHEDULED --> INTERVIEW_COMPLETED
    INTERVIEW_COMPLETED --> OFFERED
    INTERVIEW_COMPLETED --> REJECTED
    OFFERED --> ACCEPTED
    OFFERED --> DECLINED
    ACCEPTED --> ONBOARDING
    ONBOARDING --> [*]
    REJECTED --> [*]
    DECLINED --> [*]""",

    "interview_architecture": """graph TD
    subgraph Scheduling
        TRIG[Status Change] --> PANEL[Panel Assignment]
        PANEL --> CAL[Calendar Slot Selection]
        CAL --> NOTIF[Email + Calendar Invite]
        NOTIF --> CONFIRM[Candidate Confirmation]
    end

    subgraph Conduct["Interview Room"]
        JOIN[Join Room] --> AUTH[Auth Check]
        AUTH --> WEBRTC[WebRTC Session]
        WEBRTC --> VIDEO[Video Stream]
        WEBRTC --> AUDIO[Audio Stream]
        WEBRTC --> CHAT[Text Chat]
        WEBRTC --> SHARE[Screen Share]
    end

    subgraph Evaluation
        RUBRIC[Structured Rubric] --> CARD[Score Card]
        CARD --> NOTES[Interviewer Notes]
        NOTES --> SUBMIT[Submit Evaluation]
        SUBMIT --> AGG[Score Aggregation]
    end

    subgraph AI_Features["AI-Assisted"]
        TRANS[Transcription]
        SENT[Sentiment Analysis]
        SUGGEST[Question Generation]
    end

    Scheduling --> Conduct
    Conduct --> Evaluation
    Conduct -.-> AI_Features""",

    "notification_architecture": """graph LR
    subgraph Events
        E1[UserRegistered]
        E2[ApplicationSubmitted]
        E3[StatusChanged]
        E4[InterviewScheduled]
        E5[OfferGenerated]
    end

    subgraph Bus["Event Bus"]
        PUB[Publisher]
        SUB[Subscriber]
    end

    subgraph Workers
        EW[Email Worker]
        NW[In-App Worker]
    end

    subgraph Delivery
        EMAIL[SMTP/SendGrid]
        BELL[Notification Center]
    end

    Events --> PUB --> SUB
    SUB --> EW --> EMAIL
    SUB --> NW --> BELL""",

    "deployment_architecture": """graph TD
    subgraph Edge
        CDN[CDN / Firebase]
        LB[Load Balancer]
    end

    subgraph App["Application Tier"]
        API1[Flask Instance 1]
        API2[Flask Instance 2]
    end

    subgraph Worker["Worker Tier"]
        W1[Celery Email Worker]
        W2[Celery ML Worker]
        W3[Celery Beat Scheduler]
    end

    subgraph Data
        MONGO[(MongoDB Atlas)]
        REDIS[(Redis)]
        S3[File Storage]
    end

    subgraph Monitor["Observability"]
        SENTRY[Sentry]
        LOGS[Structured Logging]
        HEALTH[Health Checks]
    end

    CDN --> LB
    LB --> API1
    LB --> API2
    API1 --> REDIS
    API2 --> REDIS
    REDIS --> W1
    REDIS --> W2
    API1 --> MONGO
    API2 --> MONGO
    W1 --> LOGS
    API1 --> SENTRY""",

    "security_architecture": """graph TD
    subgraph Auth["Authentication"]
        LOCAL[Email/Password]
        GOOGLE[Google OAuth 2.0]
        JWT_TOK[JWT Token]
    end

    subgraph AuthZ["Authorization"]
        RBAC[Role-Based Access]
        PERM[Permission System]
    end

    subgraph Protection["Data Protection"]
        ENCRYPT[PII Encryption]
        ANON[Resume Anonymization]
        HASH[Password Hashing]
    end

    subgraph Defense["Defense Layer"]
        RATE[Rate Limiting]
        CORS_P[CORS Policy]
        CSP[Content Security Policy]
        XSS[XSS Prevention]
    end

    LOCAL --> JWT_TOK
    GOOGLE --> JWT_TOK
    JWT_TOK --> RBAC
    RBAC --> PERM
    ENCRYPT --> ANON
    HASH --> LOCAL
    RATE --> API[API Gateway]
    CORS_P --> API
    CSP --> API""",

    "methodology_flow": """graph TD
    REQ[Requirements Gathering] --> DES[System Design]
    DES --> PROTO[Prototype]
    PROTO --> SPRINT[Implementation Sprint]
    SPRINT --> UNIT[Unit Testing]
    UNIT --> INT_TEST[Integration Testing]
    INT_TEST --> REVIEW[Sprint Review]
    REVIEW --> FEEDBACK[Feedback Integration]
    FEEDBACK --> SPRINT
    REVIEW --> RELEASE[Release]
    RELEASE --> MONITOR[Monitor & Improve]
    MONITOR --> REQ""",

    "cicd_pipeline": """graph LR
    CODE[Git Push] --> BUILD[Docker Build]
    BUILD --> TEST[Run Tests]
    TEST --> SCAN[Security Scan]
    SCAN --> STAGE[Deploy Staging]
    STAGE --> SMOKE[Smoke Tests]
    SMOKE --> PROD[Deploy Production]
    PROD --> HEALTH[Health Check]""",

    "gantt_timeline": """gantt
    title Smart Hiring Project Timeline
    dateFormat YYYY-MM-DD
    section Phase 1 - Stabilization
    Security Critical Fixes      :crit, s1, 2026-02-16, 5d
    Architecture Consolidation   :s2, 2026-02-21, 7d
    Core Bug Fixes               :s3, 2026-02-21, 5d
    section Phase 2 - AI Enhancement
    ML Pipeline Upgrade          :s4, 2026-02-28, 14d
    Fairness Engine v2           :s5, 2026-03-07, 10d
    section Phase 3 - UI/UX
    Dashboard Redesign           :s6, 2026-03-14, 14d
    Data Visualization           :s7, 2026-03-21, 10d
    section Phase 4 - Enterprise
    Interview A/V Integration    :s8, 2026-03-28, 14d
    DevOps & Monitoring          :s9, 2026-04-04, 10d
    Security Hardening           :s10, 2026-04-11, 7d
    section Review
    MVP Demo Ready               :milestone, m1, 2026-03-01, 0d
    Feature Complete             :milestone, m2, 2026-04-15, 0d"""
}


def render_mermaid_diagrams():
    """Render all Mermaid diagrams to PNG via mermaid.ink API."""
    DIAG_DIR.mkdir(parents=True, exist_ok=True)
    paths = {}
    for key, code in MERMAID_DIAGRAMS.items():
        out = DIAG_DIR / f"{key}.png"
        if out.exists() and os.path.getsize(out) > 1000:
            paths[key] = str(out)
            print(f"    [cached] {key}")
            continue
        try:
            encoded = base64.urlsafe_b64encode(code.encode("utf-8")).decode("ascii")
            url = f"https://mermaid.ink/img/{encoded}?type=png&bgColor=white&width=1200"
            req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
            resp = urllib.request.urlopen(req, timeout=30)
            data = resp.read()
            with open(out, "wb") as f:
                f.write(data)
            img = Image.open(out)
            print(f"    [rendered] {key} ({img.width}x{img.height})")
            paths[key] = str(out)
        except KeyboardInterrupt:
            print(f"    [fallback] {key}: interrupted")
            _make_placeholder(key, out)
            paths[key] = str(out)
        except Exception as e:
            print(f"    [fallback] {key}: {e}")
            _make_placeholder(key, out)
            paths[key] = str(out)
    return paths


def _make_placeholder(title, out_path, w=1000, h=400):
    img = Image.new("RGB", (w, h), (245, 247, 250))
    d = ImageDraw.Draw(img)
    d.rectangle([3, 3, w - 4, h - 4], outline=(180, 190, 200), width=2)
    try:
        fnt = ImageFont.truetype("arial.ttf", 22)
        sfnt = ImageFont.truetype("arial.ttf", 14)
    except Exception:
        fnt = ImageFont.load_default()
        sfnt = fnt
    d.text((w // 2, h // 2 - 15), "[Architecture Diagram]", fill=(80, 80, 80), font=fnt, anchor="mm")
    d.text((w // 2, h // 2 + 25), title.replace("_", " ").title(), fill=(130, 130, 130), font=sfnt, anchor="mm")
    img.save(out_path, "PNG")


# ═══════════════════════════════════════════════════════════════
#  DOCX HELPER UTILITIES
# ═══════════════════════════════════════════════════════════════
def _shade_cell(cell, color):
    shading = parse_xml(
        f'<w:shd {nsdecls("w")} w:fill="{color[0]:02X}{color[1]:02X}{color[2]:02X}" w:val="clear"/>'
    )
    cell._tc.get_or_add_tcPr().append(shading)


def _add_table(doc, headers, rows, header_color=NAVY):
    ncols = len(headers)
    t = doc.add_table(rows=1 + len(rows), cols=ncols)
    t.style = "Table Grid"
    t.alignment = WD_TABLE_ALIGNMENT.CENTER
    # Header
    for ci, h in enumerate(headers):
        cell = t.rows[0].cells[ci]
        cell.text = h
        _shade_cell(cell, header_color)
        for run in cell.paragraphs[0].runs:
            run.font.color.rgb = RGBColor(*WHITE)
            run.bold = True
            run.font.size = Pt(9)
        cell.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
    # Data
    for ri, row in enumerate(rows):
        for ci in range(ncols):
            cell = t.rows[ri + 1].cells[ci]
            val = row[ci] if ci < len(row) else ""
            cell.text = val
            for run in cell.paragraphs[0].runs:
                run.font.size = Pt(8.5)
            if ri % 2 == 0:
                _shade_cell(cell, (240, 245, 252))
    doc.add_paragraph()
    return t


def _add_heading(doc, text, level=1):
    h = doc.add_heading(text, level=level)
    for run in h.runs:
        run.font.color.rgb = RGBColor(*NAVY)
    return h


def _add_para(doc, text, bold=False, italic=False, size=10, color=DARK_GRAY):
    p = doc.add_paragraph()
    r = p.add_run(text)
    r.bold = bold
    r.italic = italic
    r.font.size = Pt(size)
    r.font.color.rgb = RGBColor(*color)
    return p


def _add_diagram(doc, diag_paths, key, caption=""):
    path = diag_paths.get(key)
    if path and os.path.exists(path):
        try:
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            r = p.add_run()
            r.add_picture(path, width=Inches(6.0))
            if caption:
                cap = doc.add_paragraph()
                cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
                cr = cap.add_run(f"Figure: {caption}")
                cr.italic = True
                cr.font.size = Pt(9)
                cr.font.color.rgb = RGBColor(*MED_GRAY)
        except Exception:
            doc.add_paragraph(f"[Diagram: {key}]")
    else:
        doc.add_paragraph(f"[Diagram: {key}]")


# ═══════════════════════════════════════════════════════════════
#  DOCX GENERATOR — COMPREHENSIVE REVIEW DOCUMENT
# ═══════════════════════════════════════════════════════════════
def generate_review_docx(diag_paths, out_path):
    print("  Building DOCX...")
    doc = Document()

    # Styles
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(10.5)
    style.font.color.rgb = RGBColor(*DARK_GRAY)
    style.paragraph_format.space_after = Pt(4)
    style.paragraph_format.line_spacing = 1.15

    for lvl in range(1, 5):
        hs = doc.styles[f"Heading {lvl}"]
        hs.font.name = "Calibri"
        hs.font.color.rgb = RGBColor(*NAVY)
        hs.font.bold = True
        if lvl == 1: hs.font.size = Pt(22)
        elif lvl == 2: hs.font.size = Pt(16)
        elif lvl == 3: hs.font.size = Pt(13)
        else: hs.font.size = Pt(11)

    section = doc.sections[0]
    section.top_margin = Cm(2.5)
    section.bottom_margin = Cm(2)
    section.left_margin = Cm(2.5)
    section.right_margin = Cm(2.5)

    # Header/Footer
    header = section.header
    hp = header.paragraphs[0]
    hp.text = f"{PROJECT_TITLE} — Project Review Document | {DATE_STR}"
    hp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for run in hp.runs:
        run.font.size = Pt(7)
        run.font.color.rgb = RGBColor(*MED_GRAY)

    footer = section.footer
    fp = footer.paragraphs[0]
    fp.text = f"{PROJECT_TITLE} | {PROJECT_SUBTITLE} | Confidential"
    fp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for run in fp.runs:
        run.font.size = Pt(7)
        run.font.color.rgb = RGBColor(*MED_GRAY)

    # ═══════════════════════════════════════
    # TITLE PAGE
    # ═══════════════════════════════════════
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.space_before = Pt(100)
    r = p.add_run(PROJECT_TITLE.upper())
    r.bold = True; r.font.size = Pt(36); r.font.color.rgb = RGBColor(*NAVY)

    p2 = doc.add_paragraph()
    p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r2 = p2.add_run(PROJECT_SUBTITLE)
    r2.font.size = Pt(16); r2.font.color.rgb = RGBColor(*ROYAL_BLUE)

    p3 = doc.add_paragraph()
    p3.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p3.space_before = Pt(30)
    for line in [
        f"Version: {VERSION}",
        f"Date: {DATE_STR}",
        f"Team Size: {TEAM_SIZE}",
        "Institution: Academic Review Submission",
        "Classification: Project Review Document"
    ]:
        r3 = p3.add_run(line + "\n")
        r3.font.size = Pt(11)
        r3.font.color.rgb = RGBColor(*MED_GRAY)

    doc.add_page_break()

    # ═══════════════════════════════════════
    # TABLE OF CONTENTS
    # ═══════════════════════════════════════
    _add_heading(doc, "Table of Contents", 1)
    toc_items = [
        "1. Executive Summary",
        "2. Problem Statement & Motivation",
        "3. System Architecture & Design",
        "4. Technology Stack",
        "5. AI/ML Intelligence Layer",
        "6. Phase-wise Implementation",
        "   6.1 Job Posting & Sourcing",
        "   6.2 Resume Ingestion & Parsing",
        "   6.3 AI-Powered Screening",
        "   6.4 Online Assessment",
        "   6.5 Interview Phase (with Audio/Video)",
        "   6.6 Decision Engine",
        "   6.7 Offer & Onboarding",
        "7. Notification & Communication Architecture",
        "8. Fairness, Transparency & Bias Governance",
        "9. Security Architecture",
        "10. DevOps & Deployment",
        "11. Development Methodology",
        "12. Testing Strategy",
        "13. User Interface Design",
        "14. Scalability Strategy",
        "15. Risk Assessment & Mitigation",
        "16. Project Timeline & Roadmap",
        "17. Future Scope",
        "18. Conclusion",
    ]
    for item in toc_items:
        p = doc.add_paragraph(item)
        p.paragraph_format.space_after = Pt(2)
        for run in p.runs:
            run.font.size = Pt(10)
    doc.add_page_break()

    # ═══════════════════════════════════════
    # 1. EXECUTIVE SUMMARY
    # ═══════════════════════════════════════
    _add_heading(doc, "1. Executive Summary", 1)
    _add_para(doc, (
        "Smart Hiring is an AI-powered intelligent recruitment platform designed to transform "
        "the traditional hiring process by introducing data-driven decision making, bias-aware "
        "candidate evaluation, and end-to-end automation. The system addresses three critical "
        "failures in conventional recruitment: subjective screening, experience bias against "
        "freshers and career-changers, and lack of decision transparency."
    ))
    _add_para(doc, (
        "The platform implements a 9-phase hiring pipeline spanning from job posting through "
        "onboarding, with AI-powered resume parsing using spaCy NLP, semantic matching via "
        "Sentence-BERT (SBERT), fairness-constrained ranking, and explainable scoring. Every "
        "decision is audit-logged and explainable, ensuring compliance with emerging AI "
        "governance standards."
    ))
    _add_para(doc, "Key Technical Highlights:", bold=True)
    highlights = [
        "Triple-fallback ML pipeline: SBERT → TF-IDF → Keyword matching",
        "350+ skill taxonomy with context-aware extraction and anti-fraud detection",
        "Resume anonymization for bias-free evaluation (blind mode)",
        "Fresher-specific scoring: 60% skills + 25% projects + 15% certifications",
        "Event-driven notification system with async email delivery",
        "Role-based access control with 6 roles and 28 permissions",
        "Docker-containerized deployment with multi-platform support",
        "Real-time WebSocket updates for live status tracking",
    ]
    for h in highlights:
        doc.add_paragraph(h, style="List Bullet")
    doc.add_page_break()

    # ═══════════════════════════════════════
    # 2. PROBLEM STATEMENT
    # ═══════════════════════════════════════
    _add_heading(doc, "2. Problem Statement & Motivation", 1)
    _add_heading(doc, "2.1 The Problem", 2)
    _add_para(doc, (
        "Traditional recruitment suffers from multiple systemic inefficiencies that affect "
        "both organizations and candidates:"
    ))
    problems = [
        ("Subjective Screening", "78% of hiring decisions are influenced by unconscious bias. Manual resume screening is inconsistent, time-consuming, and scales poorly."),
        ("Experience Bias", "Freshers and career-changers are systematically disadvantaged by years-of-experience based filtering, ignoring skills, projects, and potential."),
        ("Opacity", "Candidates receive no feedback on why they were rejected. Scoring criteria are hidden, creating distrust in the process."),
        ("Inefficiency", "Average time-to-hire is 36 days. Manual processes create bottlenecks at every stage from screening to scheduling."),
        ("Compliance Risk", "Without audit trails, organizations face regulatory risk under emerging AI governance frameworks (EU AI Act, EEOC guidelines)."),
    ]
    for title, desc in problems:
        _add_para(doc, title, bold=True, size=11)
        _add_para(doc, desc)

    _add_heading(doc, "2.2 Our Solution", 2)
    _add_para(doc, (
        "Smart Hiring addresses these challenges through an AI-augmented recruitment pipeline "
        "that combines NLP-based resume understanding, semantic job matching, fairness-constrained "
        "ranking, and full decision transparency. The system ensures that every candidate — "
        "whether a fresher or experienced professional — is evaluated objectively based on "
        "skills, capabilities, and potential rather than superficial criteria."
    ))
    doc.add_page_break()

    # ═══════════════════════════════════════
    # 3. SYSTEM ARCHITECTURE
    # ═══════════════════════════════════════
    _add_heading(doc, "3. System Architecture & Design", 1)
    _add_heading(doc, "3.1 Architectural Approach", 2)
    _add_para(doc, (
        "The system follows a Modular Monolith architecture using Flask's Blueprint pattern "
        "for logical separation. This provides the simplicity of a monolith for a small team "
        "while maintaining the modularity needed for future microservice extraction. The "
        "architecture follows layered separation of concerns:"
    ))
    _add_table(doc,
        ["Layer", "Responsibility", "Technologies"],
        [
            ["Presentation", "User interface, client-side validation", "Vanilla JS SPA, CSS3, Chart.js"],
            ["API", "Request routing, auth, validation", "Flask Blueprints, Flask-JWT-Extended"],
            ["Service", "Business logic, ML orchestration", "Python services, SBERT, spaCy"],
            ["Data Access", "Persistence, caching", "PyMongo, Redis"],
            ["Infrastructure", "Background tasks, queuing", "Celery, Redis, Docker"],
        ]
    )

    _add_heading(doc, "3.2 High-Level System Architecture", 2)
    _add_diagram(doc, diag_paths, "system_architecture", "High-Level System Architecture")

    _add_heading(doc, "3.3 Hiring Pipeline Flow", 2)
    _add_diagram(doc, diag_paths, "hiring_pipeline", "End-to-End Hiring Pipeline")

    _add_heading(doc, "3.4 Application State Machine", 2)
    _add_para(doc, (
        "Every application follows a deterministic state machine with validated transitions. "
        "Invalid state changes are rejected, and every transition is audit-logged with the "
        "actor, timestamp, and reason."
    ))
    _add_diagram(doc, diag_paths, "state_machine", "Application Lifecycle State Machine")
    doc.add_page_break()

    # ═══════════════════════════════════════
    # 4. TECHNOLOGY STACK
    # ═══════════════════════════════════════
    _add_heading(doc, "4. Technology Stack", 1)
    _add_para(doc, (
        "The technology stack was selected to balance development velocity, ML capability, "
        "production reliability, and team expertise. Each choice is justified below."
    ))
    _add_table(doc,
        ["Category", "Technology", "Version", "Justification"],
        [
            ["Backend Framework", "Flask", "3.0.0", "Lightweight, extensible, ideal for REST APIs with ML integration"],
            ["WSGI Server", "Gunicorn", "21.2.0", "Production-grade, multi-worker, Unix-optimized"],
            ["Database", "MongoDB", "7.0", "Schema-flexible for varying resume structures, document-oriented"],
            ["Cache / Broker", "Redis", "7.x", "In-memory speed for caching, message brokering, rate limiting"],
            ["Task Queue", "Celery", "5.3.4", "Distributed task execution, retry logic, scheduling"],
            ["ML Embeddings", "Sentence-BERT", "2.2.2", "State-of-the-art semantic similarity, 384-dim embeddings"],
            ["NLP Engine", "spaCy", "3.7.2", "Fast NER, tokenization, entity recognition"],
            ["Fairness", "Fairlearn + Custom", "0.9.0", "Demographic parity, disparate impact metrics"],
            ["ML Framework", "PyTorch", "2.0+", "SBERT model inference, GPU-acceleratable"],
            ["Auth", "Flask-JWT-Extended", "4.6.0", "JWT-based stateless authentication"],
            ["OAuth", "Google OAuth 2.0", "-", "Federated identity, secure token exchange"],
            ["Real-time", "Flask-SocketIO", "5.3.5", "WebSocket support for live updates"],
            ["Document Parsing", "pdfplumber + python-docx", "-", "Multi-format resume extraction"],
            ["Security", "cryptography + bcrypt", "-", "PII encryption, password hashing"],
            ["Monitoring", "Sentry SDK", "1.39.2", "Error tracking, performance monitoring"],
            ["Containerization", "Docker", "-", "Consistent deployment, dependency isolation"],
            ["Frontend", "Vanilla JavaScript", "ES6+", "Zero-framework overhead, direct DOM control"],
        ]
    )
    doc.add_page_break()

    # ═══════════════════════════════════════
    # 5. AI/ML INTELLIGENCE LAYER
    # ═══════════════════════════════════════
    _add_heading(doc, "5. AI/ML Intelligence Layer", 1)
    _add_heading(doc, "5.1 AI Pipeline Overview", 2)
    _add_para(doc, (
        "The AI layer is the core differentiator of Smart Hiring. It replaces subjective "
        "human screening with objective, explainable, and fair AI-driven evaluation."
    ))
    _add_diagram(doc, diag_paths, "ai_pipeline", "Complete AI/ML Pipeline Architecture")

    _add_heading(doc, "5.2 Resume Parsing (NLP)", 2)
    _add_para(doc, (
        "The resume parser (1,435 lines) uses a multi-engine extraction approach with spaCy "
        "NER for entity recognition. It supports PDF and DOCX formats with cascading fallback "
        "extraction (pdfplumber → PyPDF2 → spaCy). Key capabilities:"
    ))
    parsing_features = [
        "Skill extraction from 350+ curated taxonomy with alias normalization (e.g., 'ML' → 'machine learning')",
        "Experience calculation with date merging for overlapping roles",
        "Education level detection (PhD, Masters, Bachelors, Diploma)",
        "Project and certification extraction",
        "Anti-fraud detection (keyword stuffing, hidden text, suspicious patterns)",
        "Contact information extraction (email, phone, LinkedIn)",
    ]
    for f in parsing_features:
        doc.add_paragraph(f, style="List Bullet")

    _add_heading(doc, "5.3 Semantic Matching (SBERT)", 2)
    _add_para(doc, (
        "We use Sentence-BERT with the all-MiniLM-L6-v2 model to compute dense vector "
        "representations of both resumes and job descriptions. This model produces 384-dimensional "
        "embeddings that capture semantic meaning rather than just keyword overlap."
    ))
    _add_para(doc, "The matching formula:", bold=True)
    _add_para(doc, (
        "match_score = w₁ · cos(v_resume, v_JD) + w₂ · J(S_resume, S_JD) + w₃ · f(exp) + w₄ · g(edu)"
    ), italic=True, size=11, color=ROYAL_BLUE)
    _add_para(doc, "Where:")
    formula_parts = [
        "cos(v_resume, v_JD) = Cosine similarity of SBERT embeddings (semantic match)",
        "J(S_resume, S_JD) = Jaccard index of skill sets with alias normalization",
        "f(exp) = Experience alignment function with diminishing returns above requirement",
        "g(edu) = Education level matching score",
        "w₁ + w₂ + w₃ + w₄ = 1 (configurable weights)",
    ]
    for fp in formula_parts:
        doc.add_paragraph(fp, style="List Bullet")

    _add_heading(doc, "5.4 Scoring Strategy", 2)
    _add_para(doc, (
        "The system uses differentiated scoring for freshers and experienced candidates to "
        "ensure fair evaluation regardless of career stage:"
    ))
    _add_table(doc,
        ["Component", "Fresher Weight", "Experienced Weight", "Rationale"],
        [
            ["Skills Match", "60%", "40%", "Skills are the primary indicator for freshers"],
            ["Experience", "0%", "40%", "Freshers have no work experience to evaluate"],
            ["Projects", "25%", "0%", "Academic/personal projects show fresher capability"],
            ["Certifications", "15%", "10%", "Demonstrates continuous learning"],
            ["Domain Knowledge", "0%", "10%", "Industry-specific expertise for experienced"],
        ]
    )

    _add_heading(doc, "5.5 Explainability Layer", 2)
    _add_para(doc, (
        "Every AI decision includes a human-readable explanation. The explainability service "
        "generates transparency reports showing score breakdowns, skill gap analysis, and "
        "specific match/mismatch details. This supports GDPR Article 22 right to explanation "
        "and helps recruiters understand AI recommendations."
    ))

    _add_heading(doc, "5.6 Model Inventory", 2)
    _add_table(doc,
        ["Model", "Purpose", "Size", "Status"],
        [
            ["all-MiniLM-L6-v2", "Semantic embedding (384-dim)", "80MB", "Active"],
            ["en_core_web_sm", "NER / Resume parsing", "12MB", "Active"],
            ["TfidfVectorizer", "TF-IDF fallback similarity", "5-20MB", "Active"],
            ["Custom Skill Extractor", "350+ skill taxonomy", "2MB", "Active"],
            ["Fairness Engine", "8 statistical metrics", "CPU-only", "Active"],
            ["CCI Calculator", "Career Consistency Index", "CPU-only", "Active"],
        ]
    )
    doc.add_page_break()

    # ═══════════════════════════════════════
    # 6. PHASE-WISE IMPLEMENTATION
    # ═══════════════════════════════════════
    _add_heading(doc, "6. Phase-wise Implementation", 1)
    _add_para(doc, (
        "The Smart Hiring system operates through 9 distinct phases, each with defined "
        "inputs, processing logic, ML components, database interactions, and outputs."
    ))

    # 6.1 Job Posting
    _add_heading(doc, "6.1 Job Posting & Sourcing", 2)
    _add_table(doc,
        ["Aspect", "Details"],
        [
            ["Input", "Job title, description, required skills, experience range, location"],
            ["Processing", "Skill extraction from JD using NLP, auto-tagging, normalization"],
            ["ML Component", "NER-based skill extraction, job category classification"],
            ["Database", "Jobs collection (MongoDB) — stores parsed skills, embeddings, status"],
            ["Security", "Recruiter/admin role required, input validation, XSS prevention"],
            ["Output", "Published job with extracted skills and pre-computed JD embedding"],
        ]
    )

    # 6.2 Resume Ingestion
    _add_heading(doc, "6.2 Resume Ingestion & Parsing", 2)
    _add_table(doc,
        ["Aspect", "Details"],
        [
            ["Input", "PDF/DOCX resume upload (max 5MB, validated format)"],
            ["Processing", "Multi-engine text extraction → NLP parsing → feature extraction"],
            ["ML Component", "spaCy NER, regex patterns, 350+ skill taxonomy matching"],
            ["Database", "Candidates collection — structured profile with parsed skills, experience"],
            ["Security", "File type validation, virus scanning, file size limits, path traversal prevention"],
            ["Output", "Structured candidate profile with skills, experience, education, projects"],
        ]
    )

    # 6.3 AI Screening
    _add_heading(doc, "6.3 AI-Powered Screening", 2)
    _add_table(doc,
        ["Aspect", "Details"],
        [
            ["Input", "Parsed candidate profile + job requirements"],
            ["Processing", "SBERT embedding → cosine similarity → weighted scoring → ranking"],
            ["ML Component", "SBERT all-MiniLM-L6-v2, TF-IDF fallback, weighted aggregation"],
            ["Database", "Applications collection — match_score, cci_score, skill breakdown stored"],
            ["Security", "Resume anonymization (blind mode), bias monitoring"],
            ["Output", "Ranked candidate list with explainable match scores and audit trail"],
        ]
    )

    # 6.4 Assessment
    _add_heading(doc, "6.4 Online Assessment", 2)
    _add_table(doc,
        ["Aspect", "Details"],
        [
            ["Input", "Shortlisted candidate, quiz assignment from recruiter"],
            ["Processing", "Question pool selection, timer enforcement, auto-evaluation"],
            ["ML Component", "Score normalization, adaptive difficulty (planned)"],
            ["Database", "Quizzes collection, quiz_attempts collection with per-question tracking"],
            ["Security", "Time enforcement, single-attempt validation, answer integrity checks"],
            ["Output", "Assessment score (0-100), per-question breakdown, pass/fail determination"],
        ]
    )

    # 6.5 INTERVIEW PHASE (DETAILED)
    _add_heading(doc, "6.5 Interview Phase (with Audio/Video Integration)", 2)
    _add_para(doc, (
        "The Interview Phase is the most human-intensive stage of the hiring pipeline. "
        "We design it with careful attention to scheduling, structured evaluation, and "
        "optional AI-assisted features. Audio/Video integration transforms remote interviews "
        "into professional, recorded, and analyzable sessions."
    ), bold=False)

    _add_heading(doc, "6.5.1 Interview Scheduling", 3)
    _add_para(doc, (
        "When a candidate's status transitions to INTERVIEW_SCHEDULED, the system orchestrates "
        "panel assignment, calendar slot selection, and automated notification delivery."
    ))
    _add_table(doc,
        ["Component", "Implementation"],
        [
            ["Panel Assignment", "Recruiter assigns interview panel (1-5 interviewers) with role-based selection"],
            ["Calendar Integration", "Google Calendar API or internal slot management system"],
            ["Conflict Resolution", "Smart scheduling avoids double-bookings, respects time zones"],
            ["Notification", "Email with calendar invite (.ics), interview link, preparation tips"],
            ["Confirmation", "Candidate confirms/requests reschedule via secure link"],
        ]
    )

    _add_heading(doc, "6.5.2 Audio/Video Architecture", 3)
    _add_para(doc, (
        "For the video/audio interview capability, we evaluate three technology approaches:"
    ))
    _add_table(doc,
        ["Technology", "Type", "Pros", "Cons", "Recommendation"],
        [
            ["WebRTC + LiveKit", "Open-source SFU", "Self-hosted, scalable, recording built-in, low latency", "Requires server infrastructure", "RECOMMENDED for production"],
            ["Twilio Video", "Managed SaaS", "Reliable, HIPAA-compliant, easy integration", "$0.004/min/participant, vendor lock-in", "Good for rapid development"],
            ["Jitsi Meet", "Open-source", "Free, full-featured, embeddable", "Heavy infrastructure, complex scaling", "Good for self-hosted"],
            ["Daily.co", "Managed SaaS", "Simple API, prebuilt UI components", "SaaS dependency, cost at scale", "Good for MVP"],
        ]
    )

    _add_para(doc, "Recommended Architecture: WebRTC with LiveKit", bold=True, size=11)
    _add_diagram(doc, diag_paths, "interview_architecture", "Interview Phase Architecture with A/V")

    _add_heading(doc, "6.5.3 Audio/Video Technical Requirements", 3)
    _add_table(doc,
        ["Requirement", "Priority", "Implementation"],
        [
            ["1:1 Video Call", "P0 (Critical)", "WebRTC peer connection via LiveKit SDK"],
            ["Panel Interview (3-5)", "P1 (High)", "Multi-participant SFU room with LiveKit"],
            ["Screen Sharing", "P1 (High)", "getDisplayMedia() API for technical interviews"],
            ["Text Chat", "P2 (Medium)", "WebSocket-based chat alongside video"],
            ["Recording + Consent", "P2 (Medium)", "Server-side recording with GDPR consent dialog"],
            ["Real-time Transcription", "P3 (Enhancement)", "Whisper API or Google Speech-to-Text"],
            ["Noise Suppression", "P2 (Medium)", "WebRTC built-in or Krisp SDK integration"],
            ["Bandwidth Adaptation", "P1 (High)", "Simulcast/SVC encoding, quality auto-adjustment"],
            ["Connection Quality", "P1 (High)", "Packet loss / jitter / RTT monitoring display"],
            ["Waiting Room", "P2 (Medium)", "Candidate waits until interviewer admits"],
            ["Interview Timer", "P1 (High)", "Visible countdown with configurable duration"],
            ["Mobile Support", "P1 (High)", "Responsive WebRTC with adaptive layout"],
        ]
    )

    _add_heading(doc, "6.5.4 Interview Room Implementation", 3)
    _add_para(doc, "Backend Token Generation (LiveKit):", bold=True)
    code_block = (
        "# Backend: Generate secure room token\n"
        "@interview_bp.route('/room/token', methods=['POST'])\n"
        "@jwt_required()\n"
        "def generate_room_token():\n"
        "    room_name = f'interview_{interview_id}'\n"
        "    token = livekit.AccessToken(\n"
        "        api_key=os.environ['LIVEKIT_API_KEY'],\n"
        "        api_secret=os.environ['LIVEKIT_API_SECRET']\n"
        "    )\n"
        "    token.add_grant(livekit.VideoGrant(\n"
        "        room_join=True, room=room_name,\n"
        "        can_publish=True, can_subscribe=True\n"
        "    ))\n"
        "    token.identity = current_user_id\n"
        "    return jsonify({'token': token.to_jwt(), 'room': room_name})\n"
    )
    p = doc.add_paragraph()
    r = p.add_run(code_block)
    r.font.name = "Consolas"
    r.font.size = Pt(8.5)

    _add_heading(doc, "6.5.5 Structured Evaluation Rubric", 3)
    _add_para(doc, (
        "Interviewers evaluate candidates using a structured scoring form with 5 dimensions. "
        "Each dimension is scored 1-10 with mandatory notes for justification."
    ))
    _add_table(doc,
        ["Dimension", "Weight", "Scoring Criteria"],
        [
            ["Technical Knowledge", "30%", "Depth of understanding, problem approach, accuracy"],
            ["Problem Solving", "25%", "Analytical thinking, edge case handling, optimization"],
            ["Communication", "20%", "Clarity, structuring thoughts, active listening"],
            ["Cultural Fit", "15%", "Team orientation, values alignment, adaptability"],
            ["Domain Expertise", "10%", "Industry knowledge, practical application experience"],
        ]
    )

    _add_heading(doc, "6.5.6 Interview Data Schema", 3)
    _add_para(doc, "MongoDB document structure for interview evaluation:", bold=True)
    schema_text = (
        "interview_evaluation = {\n"
        "    '_id': ObjectId,\n"
        "    'interview_id': str,\n"
        "    'application_id': str,\n"
        "    'interviewer_id': str,\n"
        "    'scores': {\n"
        "        'technical_knowledge': {'score': int, 'max': 10, 'notes': str},\n"
        "        'problem_solving': {'score': int, 'max': 10, 'notes': str},\n"
        "        'communication': {'score': int, 'max': 10, 'notes': str},\n"
        "        'cultural_fit': {'score': int, 'max': 10, 'notes': str},\n"
        "        'domain_expertise': {'score': int, 'max': 10, 'notes': str}\n"
        "    },\n"
        "    'overall_recommendation': str,  # strong_hire/hire/no_hire\n"
        "    'strengths': [str],\n"
        "    'concerns': [str],\n"
        "    'recording_url': str | None,\n"
        "    'transcript_url': str | None,\n"
        "    'duration_minutes': int,\n"
        "    'submitted_at': datetime\n"
        "}"
    )
    p = doc.add_paragraph()
    r = p.add_run(schema_text)
    r.font.name = "Consolas"
    r.font.size = Pt(8.5)

    # 6.6 Decision Engine
    _add_heading(doc, "6.6 Decision Engine", 2)
    _add_table(doc,
        ["Aspect", "Details"],
        [
            ["Input", "Resume match score + assessment score + interview evaluation"],
            ["Processing", "Weighted aggregation with configurable threshold enforcement"],
            ["ML Component", "Score normalization, threshold-based auto-shortlisting"],
            ["Override", "Recruiter can override with mandatory justification logging"],
            ["Bias Check", "Fairness engine validates decision against demographic parity"],
            ["Output", "Final decision: OFFERED, REJECTED, or WAITLISTED"],
        ]
    )

    # 6.7 Offer & Onboarding
    _add_heading(doc, "6.7 Offer Management & Onboarding", 2)
    _add_table(doc,
        ["Aspect", "Details"],
        [
            ["Input", "Approved candidate with final decision = OFFERED"],
            ["Processing", "Offer letter generation, compensation details, acceptance tracking"],
            ["Notification", "Automated email with offer details and acceptance/rejection link"],
            ["Tracking", "Real-time offer acceptance status with time-bound acceptance window"],
            ["Onboarding", "Document collection workflow, account provisioning, orientation scheduling"],
        ]
    )
    doc.add_page_break()

    # ═══════════════════════════════════════
    # 7. NOTIFICATION ARCHITECTURE
    # ═══════════════════════════════════════
    _add_heading(doc, "7. Notification & Communication Architecture", 1)
    _add_para(doc, (
        "The notification system follows an event-driven architecture where every significant "
        "action emits a domain event. Events are processed asynchronously through Celery workers "
        "to prevent API blocking."
    ))
    _add_diagram(doc, diag_paths, "notification_architecture", "Event-Driven Notification Architecture")

    _add_heading(doc, "7.1 Status Change Notification Engine", 2)
    _add_table(doc,
        ["Status Transition", "Subject Line", "Content"],
        [
            ["APPLIED", "Application Received — {job_title}", "Confirmation + next steps timeline"],
            ["SCREENING → SHORTLISTED", "Great News! You've Been Shortlisted", "Assessment instructions"],
            ["ASSESSMENT_PENDING", "Complete Your Assessment for {job_title}", "Quiz link + deadline"],
            ["ASSESSMENT_COMPLETED", "Assessment Results Available", "Score + next steps"],
            ["INTERVIEW_SCHEDULED", "Interview Scheduled — {date}", "Calendar link + preparation tips"],
            ["OFFERED", "Congratulations! Job Offer", "Offer details + acceptance link"],
            ["REJECTED", "Update Regarding Your Application", "Professional tone + encourage re-apply"],
        ]
    )

    _add_heading(doc, "7.2 Notification Governance", 2)
    _add_para(doc, (
        "The system includes debounce logic (5-minute window for rapid status changes), "
        "idempotency keys (application_id + status hash) to prevent duplicate emails, "
        "and a notification log table for delivery tracking and retry management."
    ))
    doc.add_page_break()

    # ═══════════════════════════════════════
    # 8. FAIRNESS & BIAS GOVERNANCE
    # ═══════════════════════════════════════
    _add_heading(doc, "8. Fairness, Transparency & Bias Governance", 1)
    _add_heading(doc, "8.1 Bias Mitigation Strategy", 2)
    _add_table(doc,
        ["Level", "Technique", "Implementation"],
        [
            ["Pre-processing", "Resume Anonymization", "Names, gender markers, photos removed before ML scoring"],
            ["Pre-processing", "Fresher-Specific Weights", "60% skills / 25% projects / 15% certs (no experience penalty)"],
            ["In-processing", "Skill-Weighted Scoring", "Skills dominate over years-of-experience for fairness"],
            ["Post-processing", "Demographic Parity Check", "Selection rate ≥ 80% across protected groups (4/5ths rule)"],
            ["Post-processing", "Disparate Impact Ratio", "P(positive|minority) / P(positive|majority) ≥ 0.8"],
            ["Process", "Blind Mode", "PII hidden during recruiter review to prevent confirmation bias"],
            ["Process", "Override Justification", "Mandatory text reason when recruiter overrides AI recommendation"],
            ["Audit", "Fairness Dashboard", "8 statistical metrics visualized for compliance review"],
        ]
    )

    _add_heading(doc, "8.2 Explainable AI (XAI) Compliance", 2)
    _add_para(doc, (
        "Every candidate receives a transparency report showing: score breakdown by category, "
        "matched vs missing skills, experience alignment details, and specific improvement "
        "suggestions. This supports GDPR Article 22 (right to explanation for automated "
        "decisions) and emerging AI governance regulations."
    ))
    doc.add_page_break()

    # ═══════════════════════════════════════
    # 9. SECURITY ARCHITECTURE
    # ═══════════════════════════════════════
    _add_heading(doc, "9. Security Architecture", 1)
    _add_diagram(doc, diag_paths, "security_architecture", "Security Architecture Overview")
    _add_table(doc,
        ["Control", "Implementation", "Status"],
        [
            ["Authentication", "JWT-based with bcrypt password hashing + Google OAuth 2.0", "Active"],
            ["Authorization", "RBAC with 6 roles (super_admin, admin, recruiter, candidate, analyst, auditor)", "Built"],
            ["Password Security", "bcrypt with salt, minimum complexity enforcement", "Active"],
            ["PII Encryption", "Fernet symmetric encryption with PBKDF2 key derivation", "Built"],
            ["Input Validation", "Server-side validation, XSS prevention (escapeHtml), parameterized queries", "Active"],
            ["Rate Limiting", "Token bucket algorithm, configurable per-endpoint limits", "Active"],
            ["CORS Policy", "Whitelist-based origin control", "Active"],
            ["Content Security Policy", "Script-src, style-src restrictions", "Active"],
            ["File Security", "Type validation, size limits, path traversal prevention, virus scanning", "Active"],
            ["Audit Logging", "All authentication, authorization, and data access events logged", "Active"],
        ]
    )
    doc.add_page_break()

    # ═══════════════════════════════════════
    # 10. DEVOPS & DEPLOYMENT
    # ═══════════════════════════════════════
    _add_heading(doc, "10. DevOps & Deployment", 1)
    _add_diagram(doc, diag_paths, "deployment_architecture", "Target Deployment Architecture")

    _add_heading(doc, "10.1 Containerization", 2)
    _add_para(doc, (
        "The application is containerized using Docker with multi-stage builds. "
        "docker-compose orchestrates the full stack including Flask API, MongoDB, Redis, "
        "and Celery workers. A lightweight Dockerfile.lite is available for resource-constrained "
        "environments."
    ))

    _add_heading(doc, "10.2 CI/CD Pipeline", 2)
    _add_diagram(doc, diag_paths, "cicd_pipeline", "CI/CD Pipeline Flow")

    _add_heading(doc, "10.3 Multi-Platform Deployment", 2)
    _add_table(doc,
        ["Platform", "Config File", "Use Case"],
        [
            ["Docker", "docker-compose.yml, Dockerfile", "Local development, self-hosted production"],
            ["Render", "render.yaml", "Cloud PaaS deployment (primary)"],
            ["Railway", "railway.json", "Alternative cloud deployment"],
            ["Fly.io", "fly.toml", "Edge computing deployment"],
            ["GCP Cloud Run", "cloudbuild.yaml", "Google Cloud serverless"],
            ["Firebase", "firebase.json", "Frontend static hosting"],
        ]
    )
    doc.add_page_break()

    # ═══════════════════════════════════════
    # 11. DEVELOPMENT METHODOLOGY
    # ═══════════════════════════════════════
    _add_heading(doc, "11. Development Methodology", 1)
    _add_heading(doc, "11.1 Methodology: Iterative Incremental (Hybrid Agile)", 2)
    _add_para(doc, (
        "The project follows an Iterative Incremental methodology — a hybrid approach "
        "combining Agile flexibility with milestone-driven delivery. This choice is justified by:"
    ))
    reasons = [
        "Fixed academic deadline requiring structured milestones (not pure Agile)",
        "AI/ML components needing experimental iteration (not Waterfall)",
        "Small team (4-6) unable to sustain full Scrum ceremonies",
        "Result: 2-week sprint cycles with defined deliverables per sprint",
    ]
    for r in reasons:
        doc.add_paragraph(r, style="List Bullet")

    _add_diagram(doc, diag_paths, "methodology_flow", "Development Methodology Flow")

    _add_heading(doc, "11.2 System Design Methodology", 2)
    _add_para(doc, (
        "API-First Design: All interfaces defined before implementation. "
        "Separation of Concerns: Routes → Services → Models → Database. "
        "Event-Driven: Async workflows via domain events and Celery tasks. "
        "Modular Monolith: Logical separation via Flask Blueprints with future microservice extraction path."
    ))

    _add_heading(doc, "11.3 AI/ML Methodology", 2)
    _add_table(doc,
        ["Stage", "Methodology", "Details"],
        [
            ["Data Preprocessing", "Multi-engine extraction + NLP cleaning", "pdfplumber → PyPDF2 → spaCy cascade"],
            ["Feature Engineering", "Taxonomy-based skill extraction", "350+ skills with alias normalization and context-aware matching"],
            ["Model Selection", "Transfer learning", "Pre-trained SBERT (all-MiniLM-L6-v2) — no training required"],
            ["Similarity Scoring", "Cosine similarity in embedding space", "384-dim vectors, calibrated with TF-IDF fallback"],
            ["Ranking", "Multi-feature weighted scoring", "Configurable weights per candidate type (fresher vs experienced)"],
            ["Bias Mitigation", "Pre/In/Post processing triad", "Anonymization → Skill-weighting → Demographic parity check"],
            ["Evaluation Metrics", "Precision@K, Disparate Impact Ratio", "Ranking quality + fairness measurement"],
        ]
    )
    doc.add_page_break()

    # ═══════════════════════════════════════
    # 12. TESTING STRATEGY
    # ═══════════════════════════════════════
    _add_heading(doc, "12. Testing Strategy", 1)
    _add_table(doc,
        ["Level", "Scope", "Tools", "Coverage"],
        [
            ["Unit Testing", "Scoring logic, parser functions, utility methods", "pytest", "Core algorithms"],
            ["Integration Testing", "API endpoints with test database", "Flask test client, pytest fixtures", "All routes"],
            ["System Testing", "Full hiring workflow end-to-end", "Postman collection, smoke tests", "Happy paths"],
            ["Load Testing", "Concurrent user simulation", "Locust / k6", "API throughput"],
            ["ML Validation", "Model accuracy, fairness metrics", "Custom scripts", "Scoring consistency"],
            ["Security Testing", "Auth bypass, injection, XSS", "Manual + automated", "All user inputs"],
            ["UAT", "User acceptance by stakeholders", "Demo walkthrough", "All user flows"],
        ]
    )
    doc.add_page_break()

    # ═══════════════════════════════════════
    # 13. USER INTERFACE DESIGN
    # ═══════════════════════════════════════
    _add_heading(doc, "13. User Interface Design", 1)
    _add_heading(doc, "13.1 Design Philosophy", 2)
    _add_para(doc, (
        "The UI follows a Corporate Premium design language: minimal, authoritative, "
        "data-first, with high-contrast typography and a muted color base with strong accents. "
        "The design aims to be 'professionally extraordinary' — not just clean, but "
        "confidence-inducing."
    ))

    _add_heading(doc, "13.2 Design System", 2)
    _add_table(doc,
        ["Token", "Light Mode", "Dark Mode", "Usage"],
        [
            ["Primary", "#0B1F3B (Deep Navy)", "#E3F2FD", "Headers, navigation, primary actions"],
            ["Accent", "#2563EB (Royal Blue)", "#60A5FA", "CTAs, links, active states"],
            ["Surface", "#FFFFFF", "#1E293B", "Card backgrounds, panels"],
            ["Text", "#1E293B", "#F1F5F9", "Body text, labels"],
            ["Success", "#10B981 (Emerald)", "#34D399", "Positive actions, confirmations"],
            ["Warning", "#F59E0B (Amber)", "#FBBF24", "Caution states, pending items"],
            ["Danger", "#EF4444 (Red)", "#F87171", "Errors, critical alerts, rejections"],
        ]
    )

    _add_heading(doc, "13.3 Dashboard Components", 2)
    dash_components = [
        "KPI Cards: Total Applications, Shortlisted, Interviews, Offers, Avg Match Score",
        "Hiring Funnel Visualization: Application → Screening → Assessment → Interview → Offer",
        "Score Distribution Chart: Histogram of match scores per job",
        "Real-time Activity Feed: Latest actions across the system",
        "System Health Indicators: Worker status, queue depth, API response times",
    ]
    for dc in dash_components:
        doc.add_paragraph(dc, style="List Bullet")

    _add_heading(doc, "13.4 Accessibility (WCAG 2.1)", 2)
    _add_para(doc, (
        "The system includes an accessibility layer (a11y.js + a11y.css) supporting "
        "keyboard navigation, screen reader compatibility, focus indicators, and "
        "sufficient color contrast ratios (minimum 4.5:1)."
    ))
    doc.add_page_break()

    # ═══════════════════════════════════════
    # 14. SCALABILITY STRATEGY
    # ═══════════════════════════════════════
    _add_heading(doc, "14. Scalability Strategy", 1)
    _add_table(doc,
        ["Strategy", "Implementation", "Impact"],
        [
            ["Horizontal Scaling", "Docker replicas behind load balancer", "Handle concurrent users"],
            ["Async Processing", "Celery workers for ML, email, reports", "Non-blocking API responses"],
            ["Caching", "Redis for session data, computed scores", "Reduced DB load, faster responses"],
            ["Database Indexing", "Compound indexes on frequent queries", "Query performance optimization"],
            ["Model Serving Isolation", "ML service can run as separate container", "Independent scaling of ML workload"],
            ["Stateless Design", "JWT auth, no server-side sessions", "Enable horizontal scaling"],
            ["CDN for Static Assets", "Firebase/Cloudflare for frontend", "Reduced origin server load"],
        ]
    )
    doc.add_page_break()

    # ═══════════════════════════════════════
    # 15. RISK ASSESSMENT
    # ═══════════════════════════════════════
    _add_heading(doc, "15. Risk Assessment & Mitigation", 1)
    _add_table(doc,
        ["Risk", "Likelihood", "Impact", "Mitigation"],
        [
            ["SBERT model fails to load", "Low", "High", "TF-IDF + keyword fallback chain ensures scoring continues"],
            ["Database outage", "Low", "Critical", "MongoDB Atlas replication, connection retry logic"],
            ["Email delivery failure", "Medium", "Medium", "Celery retry with exponential backoff, DLQ handler"],
            ["OAuth misconfiguration", "Medium", "High", "Environment-based config, fallback to email/password auth"],
            ["Scoring inconsistency", "Medium", "High", "Unified scoring weights config, regression test suite"],
            ["Resume parsing crash", "Low", "Medium", "Multi-engine fallback, graceful error handling"],
            ["Interview A/V failure", "Medium", "High", "Fallback to phone interview, connection quality monitoring"],
            ["Model bias / drift", "Medium", "Critical", "Fairness monitoring dashboard, regular bias audits"],
        ]
    )
    doc.add_page_break()

    # ═══════════════════════════════════════
    # 16. PROJECT TIMELINE & ROADMAP
    # ═══════════════════════════════════════
    _add_heading(doc, "16. Project Timeline & Roadmap", 1)
    _add_diagram(doc, diag_paths, "gantt_timeline", "Project Timeline (Gantt Chart)")

    _add_heading(doc, "16.1 Milestone Dates", 2)
    _add_table(doc,
        ["Milestone", "Target Date", "Confidence", "Deliverables"],
        [
            ["Review Demo Ready", "March 1, 2026", "85%", "Core flow functional, PPTX/DOCX ready"],
            ["MVP Complete", "March 25, 2026", "75%", "All phases operational, security hardened"],
            ["Feature Complete", "April 15, 2026", "70%", "Interview A/V, exports, full analytics"],
            ["Production Stable", "May 6, 2026", "65%", "Load tested, pen tested, monitored"],
            ["Enterprise Optimized", "May 27, 2026", "55%", "Fully scalable, governed, documented"],
        ]
    )

    _add_heading(doc, "16.2 Phase Breakdown", 2)
    _add_table(doc,
        ["Phase", "Duration", "Key Deliverables"],
        [
            ["Phase 1: Stabilization", "2 weeks", "Critical bug fixes, security hardening, architecture consolidation"],
            ["Phase 2: AI Enhancement", "3 weeks", "ML pipeline upgrade, FAISS vector store, real fairness data"],
            ["Phase 3: UI/UX Elevation", "3 weeks", "Component library, charts, responsive design, data visualization"],
            ["Phase 4: Interview + A/V", "3 weeks", "LiveKit integration, evaluation rubric, scheduling system"],
            ["Phase 5: DevOps", "2 weeks", "CI/CD pipeline, staging environment, monitoring dashboard"],
            ["Phase 6: Governance", "2 weeks", "Audit logging, GDPR compliance, model versioning, pen testing"],
        ]
    )
    doc.add_page_break()

    # ═══════════════════════════════════════
    # 17. FUTURE SCOPE
    # ═══════════════════════════════════════
    _add_heading(doc, "17. Future Scope", 1)
    future_items = [
        ("Cross-Encoder Reranking", "Replace single-stage ranking with SBERT cross-encoder for 15-25% accuracy improvement"),
        ("Skill Graph Knowledge Base", "Neo4j-based skill ontology for semantic skill relationship understanding"),
        ("LLM-Based Summary Generation", "GPT/T5-based candidate summary for recruiter quick review"),
        ("Interview Prediction Model", "XGBoost model predicting interview success from assessment + resume features"),
        ("Multi-Tenant SaaS", "Organization isolation, custom branding, per-tenant model configuration"),
        ("Internal Mobility Engine", "Match existing employees to new internal roles for talent retention"),
        ("Recruiter AI Copilot", "AI assistant for JD writing, candidate shortlisting suggestions, analytics insights"),
        ("Mobile Application", "React Native cross-platform app for candidates and recruiters"),
    ]
    for title, desc in future_items:
        _add_para(doc, title, bold=True, size=11)
        _add_para(doc, desc)
    doc.add_page_break()

    # ═══════════════════════════════════════
    # 18. CONCLUSION
    # ═══════════════════════════════════════
    _add_heading(doc, "18. Conclusion", 1)
    _add_para(doc, (
        "Smart Hiring represents a comprehensive approach to modernizing recruitment through "
        "AI-powered intelligence, ethical governance, and production-grade engineering. The "
        "system addresses the fundamental challenges of subjective screening, experience bias, "
        "and decision opacity through a carefully designed pipeline that combines NLP, semantic "
        "matching, fairness constraints, and explainable AI."
    ))
    _add_para(doc, (
        "The architecture is designed for production scalability from day one, with Docker "
        "containerization, async background processing, and multi-platform deployment support. "
        "The security model includes JWT authentication, RBAC authorization, PII encryption, "
        "and comprehensive audit logging."
    ))
    _add_para(doc, (
        "With the completed AI pipeline, assessment engine, and notification infrastructure, "
        "and with the interview A/V integration, export system, and enterprise monitoring "
        "on the near-term roadmap, Smart Hiring is positioned to become a complete, "
        "production-ready intelligent recruitment platform that brings transparency, fairness, "
        "and efficiency to every stage of the hiring process."
    ))

    doc.save(str(out_path))
    print(f"  DOCX saved: {out_path} ({os.path.getsize(out_path) // 1024} KB)")


# ═══════════════════════════════════════════════════════════════
#  PPTX GENERATOR — REVIEW PRESENTATION
# ═══════════════════════════════════════════════════════════════
def generate_review_pptx(diag_paths, out_path):
    print("  Building PPTX...")
    prs = Presentation()
    prs.slide_width = PI(13.333)
    prs.slide_height = PI(7.5)

    def add_slide():
        return prs.slides.add_slide(prs.slide_layouts[6])

    def tb(slide, l, t, w, h):
        return slide.shapes.add_textbox(PI(l), PI(t), PI(w), PI(h))

    def bg(slide, color):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = PC(*color)

    def rect(slide, l, t, w, h, color):
        shape = slide.shapes.add_shape(1, PI(l), PI(t), PI(w), PI(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = PC(*color)
        shape.line.fill.background()
        return shape

    def header_bar(slide, text, color=NAVY):
        rect(slide, 0, 0, 13.333, 0.85, color)
        box = tb(slide, 0.5, 0.1, 12, 0.65)
        p = box.text_frame.paragraphs[0]
        p.text = text
        p.font.size = PP(28)
        p.font.bold = True
        p.font.color.rgb = PC(*WHITE)

    def bullet_slide(slide, items, start_y=1.2, font_size=14):
        y = start_y
        for item in items:
            rect(slide, 1, y, 0.12, 0.5, ROYAL_BLUE)
            box = tb(slide, 1.5, y, 11, 0.5)
            p = box.text_frame.paragraphs[0]
            p.text = item
            p.font.size = PP(font_size)
            p.font.color.rgb = PC(*DARK_GRAY)
            y += 0.62
        return y

    def add_diagram_slide(title, key, color=DARK_BLUE):
        path = diag_paths.get(key)
        if not path or not os.path.exists(path):
            return
        s = add_slide()
        bg(s, WHITE)
        header_bar(s, title, color)
        try:
            img = Image.open(path)
            iw, ih = img.size
            max_w, max_h = 11.5, 5.8
            ratio = min(max_w / (iw / 96), max_h / (ih / 96))
            w = iw / 96 * ratio
            h = ih / 96 * ratio
            left = (13.333 - w) / 2
            top = 1.0 + (5.8 - h) / 2
            s.shapes.add_picture(path, PI(left), PI(top), PI(w), PI(h))
        except Exception:
            pass

    # ═══════════════════════════════
    # SLIDE 1: Title
    # ═══════════════════════════════
    s = add_slide()
    bg(s, NAVY)
    rect(s, 0, 3.5, 13.333, 0.06, STEEL)

    box = tb(s, 1, 0.8, 11, 2.5)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Smart Hiring"
    p.font.size = PP(48)
    p.font.bold = True
    p.font.color.rgb = PC(*WHITE)
    p = tf.add_paragraph()
    p.text = "AI-Powered Intelligent Recruitment Platform"
    p.font.size = PP(24)
    p.font.color.rgb = PC(*STEEL)

    box2 = tb(s, 1, 4.0, 11, 2.5)
    tf2 = box2.text_frame
    tf2.word_wrap = True
    for txt in [
        f"Version: {VERSION}  |  Date: {DATE_STR}",
        f"Team: {TEAM_SIZE}",
        "Project Review Presentation",
        "Technologies: Flask • MongoDB • SBERT • spaCy • Redis • Celery • Docker"
    ]:
        p = tf2.add_paragraph()
        p.text = txt
        p.font.size = PP(14)
        p.font.color.rgb = PC(180, 200, 220)
        p.space_after = PP(6)

    # ═══════════════════════════════
    # SLIDE 2: Why This Matters
    # ═══════════════════════════════
    s = add_slide()
    bg(s, WHITE)
    header_bar(s, "THE PROBLEM WE SOLVE", RED_ACCENT)

    problems_data = [
        ("78%", "of hiring decisions influenced by unconscious bias"),
        ("36 days", "average time-to-hire in traditional recruitment"),
        ("0%", "feedback given to rejected candidates in most systems"),
        ("Freshers", "systematically disadvantaged by experience-first filtering"),
    ]
    y = 1.3
    for stat, desc in problems_data:
        rect(s, 1, y, 2.2, 0.8, NAVY)
        box = tb(s, 1.1, y + 0.1, 2, 0.6)
        p = box.text_frame.paragraphs[0]
        p.text = stat
        p.font.size = PP(22)
        p.font.bold = True
        p.font.color.rgb = PC(*WHITE)
        p.alignment = PP_ALIGN.CENTER

        box2 = tb(s, 3.5, y + 0.15, 9, 0.5)
        p2 = box2.text_frame.paragraphs[0]
        p2.text = desc
        p2.font.size = PP(16)
        p2.font.color.rgb = PC(*DARK_GRAY)
        y += 1.1

    box3 = tb(s, 1, y + 0.5, 11.333, 0.8)
    tf3 = box3.text_frame
    tf3.word_wrap = True
    p = tf3.paragraphs[0]
    p.text = "Smart Hiring replaces subjective screening with AI-driven, bias-aware, explainable recruitment intelligence."
    p.font.size = PP(16)
    p.font.bold = True
    p.font.color.rgb = PC(*ROYAL_BLUE)
    p.alignment = PP_ALIGN.CENTER

    # ═══════════════════════════════
    # SLIDE 3: Solution Overview
    # ═══════════════════════════════
    s = add_slide()
    bg(s, WHITE)
    header_bar(s, "OUR SOLUTION")

    solutions = [
        "AI-powered resume parsing with NLP (spaCy + 350+ skill taxonomy)",
        "Semantic job matching using Sentence-BERT embeddings",
        "Fairness-constrained ranking with bias detection",
        "Fresher-specific scoring eliminating experience bias",
        "Explainable AI with GDPR Article 22 transparency reports",
        "End-to-end automation: Job Posting → Screening → Assessment → Interview → Offer",
        "Event-driven notifications with async email delivery",
        "Production-ready: Docker + CI/CD + monitoring",
    ]
    bullet_slide(s, solutions, 1.1, 14)

    # ═══════════════════════════════
    # SLIDE 4-5: Architecture Diagrams
    # ═══════════════════════════════
    add_diagram_slide("SYSTEM ARCHITECTURE", "system_architecture")
    add_diagram_slide("HIRING PIPELINE FLOW", "hiring_pipeline")

    # ═══════════════════════════════
    # SLIDE 6: Technology Stack
    # ═══════════════════════════════
    s = add_slide()
    bg(s, WHITE)
    header_bar(s, "TECHNOLOGY STACK")

    stack = [
        ("Backend", "Flask 3.0 + Gunicorn"),
        ("Database", "MongoDB 7.0"),
        ("Cache/Queue", "Redis 7 + Celery 5.3"),
        ("ML Engine", "SBERT all-MiniLM-L6-v2"),
        ("NLP", "spaCy en_core_web_sm"),
        ("Fairness", "Custom Engine + Fairlearn"),
        ("Auth", "JWT + Google OAuth 2.0"),
        ("Real-time", "Flask-SocketIO (WebSocket)"),
        ("Deploy", "Docker + docker-compose"),
        ("Monitor", "Sentry + Structured Logging"),
    ]
    y = 1.1
    for label, val in stack:
        rect(s, 1, y, 3.2, 0.52, (235, 240, 248))
        box = tb(s, 1.2, y + 0.03, 3, 0.46)
        p = box.text_frame.paragraphs[0]
        p.text = label
        p.font.size = PP(12)
        p.font.bold = True
        p.font.color.rgb = PC(*NAVY)
        box2 = tb(s, 4.4, y + 0.03, 8, 0.46)
        p2 = box2.text_frame.paragraphs[0]
        p2.text = val
        p2.font.size = PP(12)
        p2.font.color.rgb = PC(*DARK_GRAY)
        y += 0.58

    # ═══════════════════════════════
    # SLIDE 7: AI Pipeline
    # ═══════════════════════════════
    add_diagram_slide("AI/ML INTELLIGENCE PIPELINE", "ai_pipeline", ROYAL_BLUE)

    # ═══════════════════════════════
    # SLIDE 8: Scoring Strategy
    # ═══════════════════════════════
    s = add_slide()
    bg(s, WHITE)
    header_bar(s, "SCORING STRATEGY — FRESHER vs EXPERIENCED")

    # Fresher column
    rect(s, 0.5, 1.2, 5.8, 0.6, EMERALD)
    box = tb(s, 0.7, 1.25, 5.4, 0.5)
    p = box.text_frame.paragraphs[0]
    p.text = "FRESHER SCORING"
    p.font.size = PP(18)
    p.font.bold = True
    p.font.color.rgb = PC(*WHITE)
    p.alignment = PP_ALIGN.CENTER

    fresher_scores = [("Skills Match", "60%"), ("Projects", "25%"), ("Certifications", "15%"), ("Experience", "0% (no penalty)")]
    y = 2.0
    for name, pct in fresher_scores:
        rect(s, 0.7, y, 5.4, 0.5, (240, 252, 240))
        box = tb(s, 0.9, y + 0.05, 3.5, 0.4)
        p = box.text_frame.paragraphs[0]
        p.text = name
        p.font.size = PP(13)
        p.font.color.rgb = PC(*DARK_GRAY)
        box2 = tb(s, 4.5, y + 0.05, 1.5, 0.4)
        p2 = box2.text_frame.paragraphs[0]
        p2.text = pct
        p2.font.size = PP(14)
        p2.font.bold = True
        p2.font.color.rgb = PC(*EMERALD)
        p2.alignment = PP_ALIGN.RIGHT
        y += 0.55

    # Experienced column
    rect(s, 7, 1.2, 5.8, 0.6, ROYAL_BLUE)
    box = tb(s, 7.2, 1.25, 5.4, 0.5)
    p = box.text_frame.paragraphs[0]
    p.text = "EXPERIENCED SCORING"
    p.font.size = PP(18)
    p.font.bold = True
    p.font.color.rgb = PC(*WHITE)
    p.alignment = PP_ALIGN.CENTER

    exp_scores = [("Skills Match", "40%"), ("Experience", "40%"), ("Domain Knowledge", "10%"), ("Certifications", "10%")]
    y = 2.0
    for name, pct in exp_scores:
        rect(s, 7.2, y, 5.4, 0.5, (235, 240, 252))
        box = tb(s, 7.4, y + 0.05, 3.5, 0.4)
        p = box.text_frame.paragraphs[0]
        p.text = name
        p.font.size = PP(13)
        p.font.color.rgb = PC(*DARK_GRAY)
        box2 = tb(s, 11.0, y + 0.05, 1.5, 0.4)
        p2 = box2.text_frame.paragraphs[0]
        p2.text = pct
        p2.font.size = PP(14)
        p2.font.bold = True
        p2.font.color.rgb = PC(*ROYAL_BLUE)
        p2.alignment = PP_ALIGN.RIGHT
        y += 0.55

    # Key insight
    box = tb(s, 1, 5.0, 11.333, 1.0)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Key Innovation: Freshers are never penalized for lack of experience. Skills and projects receive proportionally higher weight to ensure fair evaluation."
    p.font.size = PP(13)
    p.font.italic = True
    p.font.color.rgb = PC(*MED_GRAY)
    p.alignment = PP_ALIGN.CENTER

    # ═══════════════════════════════
    # SLIDE 9: Fairness & Bias
    # ═══════════════════════════════
    s = add_slide()
    bg(s, WHITE)
    header_bar(s, "FAIRNESS & BIAS GOVERNANCE", EMERALD)

    fairness_items = [
        "Pre-processing: Resume anonymization removes names, gender markers, photos",
        "In-processing: Skill-weighted scoring reduces experience dominance",
        "Post-processing: Demographic parity + disparate impact monitoring",
        "Blind Mode: PII hidden during recruiter review to prevent bias",
        "Override Logging: Mandatory justification when overriding AI decision",
        "Explainable AI: Per-candidate transparency report with score breakdown",
        "Fairness Dashboard: 8 statistical metrics for compliance audit",
        "GDPR Art. 22: Right to explanation for automated decisions",
    ]
    bullet_slide(s, fairness_items, 1.1, 13)

    # ═══════════════════════════════
    # SLIDE 10: Interview Phase
    # ═══════════════════════════════
    add_diagram_slide("INTERVIEW PHASE ARCHITECTURE (with A/V)", "interview_architecture", ROYAL_BLUE)

    # ═══════════════════════════════
    # SLIDE 11: Interview A/V Details
    # ═══════════════════════════════
    s = add_slide()
    bg(s, WHITE)
    header_bar(s, "INTERVIEW A/V REQUIREMENTS")

    av_items = [
        "WebRTC + LiveKit (open-source SFU) — self-hosted, scalable",
        "1:1 and panel interviews (3-5 participants)",
        "Screen sharing for technical interviews",
        "Server-side recording with GDPR consent dialog",
        "Real-time transcription via Whisper / Google STT",
        "Bandwidth adaptation with simulcast encoding",
        "Structured evaluation rubric with 5-dimension scoring",
        "Interview timer, waiting room, connection quality display",
    ]
    bullet_slide(s, av_items, 1.1, 13)

    # ═══════════════════════════════
    # SLIDE 12: Security
    # ═══════════════════════════════
    add_diagram_slide("SECURITY ARCHITECTURE", "security_architecture", RED_ACCENT)

    # ═══════════════════════════════
    # SLIDE 13: Notification System
    # ═══════════════════════════════
    add_diagram_slide("NOTIFICATION ARCHITECTURE", "notification_architecture")

    # ═══════════════════════════════
    # SLIDE 14: State Machine
    # ═══════════════════════════════
    add_diagram_slide("APPLICATION STATE MACHINE", "state_machine")

    # ═══════════════════════════════
    # SLIDE 15: Deployment
    # ═══════════════════════════════
    add_diagram_slide("DEPLOYMENT ARCHITECTURE", "deployment_architecture")

    # ═══════════════════════════════
    # SLIDE 16: Methodology
    # ═══════════════════════════════
    s = add_slide()
    bg(s, WHITE)
    header_bar(s, "DEVELOPMENT METHODOLOGY")

    method_items = [
        "Iterative Incremental (Hybrid Agile) — 2-week sprint cycles",
        "API-First Design: Interfaces defined before implementation",
        "Modular Monolith: Flask Blueprints for logical separation",
        "Event-Driven Architecture: Domain events + async workers",
        "Transfer Learning: Pre-trained SBERT — no custom training required",
        "Bias Mitigation Triad: Pre/In/Post processing pipeline",
        "Test Pyramid: Unit → Integration → System → Load → ML Validation",
        "CI/CD: Docker build → Test → Scan → Stage → Deploy",
    ]
    bullet_slide(s, method_items, 1.1, 13)

    # ═══════════════════════════════
    # SLIDE 17: CI/CD
    # ═══════════════════════════════
    add_diagram_slide("CI/CD PIPELINE", "cicd_pipeline")

    # ═══════════════════════════════
    # SLIDE 18: Timeline
    # ═══════════════════════════════
    add_diagram_slide("PROJECT TIMELINE", "gantt_timeline")

    # ═══════════════════════════════
    # SLIDE 19: Milestones
    # ═══════════════════════════════
    s = add_slide()
    bg(s, WHITE)
    header_bar(s, "PROJECT MILESTONES")

    milestones = [
        ("Review Demo Ready", "Mar 1, 2026", "85%", EMERALD),
        ("MVP Complete", "Mar 25, 2026", "75%", EMERALD),
        ("Feature Complete", "Apr 15, 2026", "70%", AMBER),
        ("Production Stable", "May 6, 2026", "65%", AMBER),
        ("Enterprise Optimized", "May 27, 2026", "55%", STEEL),
    ]
    y = 1.2
    for name, date, conf, color in milestones:
        rect(s, 1, y, 0.15, 0.7, color)
        rect(s, 1.3, y, 11, 0.7, (245, 248, 252))
        box = tb(s, 1.5, y + 0.08, 5, 0.54)
        p = box.text_frame.paragraphs[0]
        p.text = name
        p.font.size = PP(16)
        p.font.bold = True
        p.font.color.rgb = PC(*DARK_GRAY)
        box2 = tb(s, 7, y + 0.08, 3, 0.54)
        p2 = box2.text_frame.paragraphs[0]
        p2.text = date
        p2.font.size = PP(14)
        p2.font.color.rgb = PC(*MED_GRAY)
        box3 = tb(s, 10, y + 0.08, 2, 0.54)
        p3 = box3.text_frame.paragraphs[0]
        p3.text = f"Conf: {conf}"
        p3.font.size = PP(13)
        p3.font.bold = True
        p3.font.color.rgb = PC(*color)
        y += 0.85

    # ═══════════════════════════════
    # SLIDE 20: Future Scope
    # ═══════════════════════════════
    s = add_slide()
    bg(s, WHITE)
    header_bar(s, "FUTURE SCOPE", DARK_BLUE)

    future = [
        "Cross-Encoder Reranking — 15-25% ranking quality improvement",
        "Skill Graph Knowledge Base (Neo4j / NetworkX)",
        "LLM-Based Candidate Summary Generation (T5/GPT)",
        "Interview Prediction Model (XGBoost)",
        "Multi-Tenant SaaS Architecture",
        "Internal Mobility Engine for employee role matching",
        "Recruiter AI Copilot for intelligent assistance",
        "Mobile Application (React Native)",
    ]
    bullet_slide(s, future, 1.1, 14)

    # ═══════════════════════════════
    # SLIDE 21: Demo Strategy
    # ═══════════════════════════════
    s = add_slide()
    bg(s, WHITE)
    header_bar(s, "LIVE DEMO FLOW")

    demo_steps = [
        "1. Recruiter creates a job posting → Show skill extraction",
        "2. Candidate registers → Resume upload and NLP parsing",
        "3. AI Screening → Show match scores and explainability report",
        "4. Assessment → Candidate takes quiz, auto-evaluation",
        "5. Dashboard → Show analytics, score distributions, fairness metrics",
        "6. Blind Mode → Toggle anonymized view for bias-free review",
        "7. Export Report → CSV download of hiring analytics",
    ]
    bullet_slide(s, demo_steps, 1.1, 14)

    # ═══════════════════════════════
    # SLIDE 22: Thank You
    # ═══════════════════════════════
    s = add_slide()
    bg(s, NAVY)
    rect(s, 0, 3.3, 13.333, 0.06, STEEL)

    box = tb(s, 2, 1.0, 9, 2.5)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = PP(48)
    p.font.bold = True
    p.font.color.rgb = PC(*WHITE)
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "Questions & Discussion"
    p.font.size = PP(24)
    p.font.color.rgb = PC(*STEEL)
    p.alignment = PP_ALIGN.CENTER

    box2 = tb(s, 2, 4.0, 9, 2.5)
    tf2 = box2.text_frame
    tf2.word_wrap = True
    for txt in [
        f"{PROJECT_TITLE} — {PROJECT_SUBTITLE}",
        f"Version {VERSION} | {DATE_STR}",
        "Flask • MongoDB • SBERT • spaCy • Redis • Celery • Docker",
        "AI-Driven • Bias-Aware • Explainable • Production-Ready"
    ]:
        p = tf2.add_paragraph()
        p.text = txt
        p.font.size = PP(13)
        p.font.color.rgb = PC(150, 170, 190)
        p.alignment = PP_ALIGN.CENTER

    prs.save(str(out_path))
    print(f"  PPTX saved: {out_path} ({os.path.getsize(out_path) // 1024} KB, {len(prs.slides)} slides)")


# ═══════════════════════════════════════════════════════════════
#  PDF GENERATOR (using reportlab)
# ═══════════════════════════════════════════════════════════════
def generate_review_pdf(diag_paths, out_path):
    if not HAS_REPORTLAB:
        print("  [SKIP] PDF generation — reportlab not installed")
        return

    print("  Building PDF...")
    doc = SimpleDocTemplate(
        str(out_path), pagesize=A4,
        topMargin=2.5*cm, bottomMargin=2*cm,
        leftMargin=2.5*cm, rightMargin=2.5*cm,
        title=f"{PROJECT_TITLE} — Project Review",
        author="Smart Hiring Team"
    )

    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(
        'CustomTitle', parent=styles['Title'],
        fontSize=28, textColor=HexColor('#0B1F3B'),
        spaceAfter=20, alignment=TA_CENTER
    ))
    styles.add(ParagraphStyle(
        'CustomH1', parent=styles['Heading1'],
        fontSize=18, textColor=HexColor('#0B1F3B'),
        spaceAfter=10, spaceBefore=20
    ))
    styles.add(ParagraphStyle(
        'CustomH2', parent=styles['Heading2'],
        fontSize=14, textColor=HexColor('#003366'),
        spaceAfter=8, spaceBefore=14
    ))
    styles.add(ParagraphStyle(
        'CustomBody', parent=styles['Normal'],
        fontSize=10, textColor=HexColor('#1E293B'),
        spaceAfter=6, leading=14, alignment=TA_JUSTIFY
    ))
    styles.add(ParagraphStyle(
        'CustomBullet', parent=styles['Normal'],
        fontSize=10, textColor=HexColor('#1E293B'),
        spaceAfter=4, leftIndent=20, bulletIndent=10,
        bulletFontName='Helvetica', bulletFontSize=10
    ))

    story = []

    # Title page
    story.append(Spacer(1, 3*cm))
    story.append(Paragraph(PROJECT_TITLE.upper(), styles['CustomTitle']))
    story.append(Paragraph(PROJECT_SUBTITLE, ParagraphStyle(
        'Subtitle', parent=styles['Normal'],
        fontSize=16, textColor=HexColor('#2563EB'),
        alignment=TA_CENTER, spaceAfter=30
    )))
    story.append(Spacer(1, 1*cm))
    meta_text = (
        f"Version: {VERSION}<br/>"
        f"Date: {DATE_STR}<br/>"
        f"Team: {TEAM_SIZE}<br/>"
        "Classification: Project Review Document"
    )
    story.append(Paragraph(meta_text, ParagraphStyle(
        'Meta', parent=styles['Normal'],
        fontSize=11, textColor=HexColor('#64748B'),
        alignment=TA_CENTER, leading=16
    )))
    story.append(PageBreak())

    # Executive Summary
    story.append(Paragraph("1. Executive Summary", styles['CustomH1']))
    story.append(Paragraph(
        "Smart Hiring is an AI-powered intelligent recruitment platform designed to transform "
        "the traditional hiring process through data-driven decision making, bias-aware "
        "candidate evaluation, and end-to-end automation. The system implements a 9-phase "
        "hiring pipeline with AI-powered resume parsing (spaCy NLP), semantic matching "
        "(Sentence-BERT), fairness-constrained ranking, and explainable scoring.",
        styles['CustomBody']
    ))
    story.append(Paragraph(
        "Key highlights: Triple-fallback ML pipeline, 350+ skill taxonomy, resume "
        "anonymization for bias-free evaluation, fresher-specific scoring that eliminates "
        "experience penalty, and event-driven notification system with async delivery.",
        styles['CustomBody']
    ))

    # Technology Stack
    story.append(Paragraph("2. Technology Stack", styles['CustomH1']))
    tech_data = [
        ['Category', 'Technology', 'Purpose'],
        ['Backend', 'Flask 3.0', 'REST API framework'],
        ['Database', 'MongoDB 7.0', 'Document-oriented persistence'],
        ['Cache/Queue', 'Redis + Celery', 'Caching, message broker, task queue'],
        ['ML Engine', 'SBERT all-MiniLM-L6-v2', 'Semantic embedding (384-dim)'],
        ['NLP', 'spaCy 3.7', 'Named Entity Recognition, parsing'],
        ['Fairness', 'Custom + Fairlearn', 'Bias detection & mitigation'],
        ['Auth', 'JWT + OAuth 2.0', 'Authentication & authorization'],
        ['Deploy', 'Docker', 'Containerization & orchestration'],
    ]
    t = Table(tech_data, colWidths=[3*cm, 4*cm, 8*cm])
    t.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), HexColor('#0B1F3B')),
        ('TEXTCOLOR', (0, 0), (-1, 0), HexColor('#FFFFFF')),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, -1), 9),
        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
        ('GRID', (0, 0), (-1, -1), 0.5, HexColor('#CCCCCC')),
        ('ROWBACKGROUNDS', (0, 1), (-1, -1), [HexColor('#FFFFFF'), HexColor('#F0F5FC')]),
        ('TOPPADDING', (0, 0), (-1, -1), 4),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
    ]))
    story.append(t)
    story.append(Spacer(1, 0.5*cm))

    # Add diagrams
    for key, caption in [
        ("system_architecture", "System Architecture"),
        ("ai_pipeline", "AI/ML Pipeline"),
        ("hiring_pipeline", "Hiring Pipeline Flow"),
        ("state_machine", "Application State Machine"),
        ("interview_architecture", "Interview Phase Architecture"),
        ("notification_architecture", "Notification Architecture"),
        ("deployment_architecture", "Deployment Architecture"),
        ("security_architecture", "Security Architecture"),
        ("methodology_flow", "Development Methodology"),
        ("cicd_pipeline", "CI/CD Pipeline"),
    ]:
        path = diag_paths.get(key)
        if path and os.path.exists(path):
            story.append(Paragraph(f"Figure: {caption}", styles['CustomH2']))
            try:
                img = RLImage(path, width=15*cm, height=8*cm, kind='proportional')
                story.append(img)
            except Exception:
                story.append(Paragraph(f"[Diagram: {key}]", styles['CustomBody']))
            story.append(Spacer(1, 0.5*cm))

    # Scoring Strategy
    story.append(PageBreak())
    story.append(Paragraph("3. Scoring Strategy", styles['CustomH1']))
    score_data = [
        ['Component', 'Fresher Weight', 'Experienced Weight', 'Rationale'],
        ['Skills Match', '60%', '40%', 'Primary capability indicator'],
        ['Experience', '0%', '40%', 'Not applicable for freshers'],
        ['Projects', '25%', '0%', 'Demonstrates fresher capability'],
        ['Certifications', '15%', '10%', 'Continuous learning indicator'],
        ['Domain Knowledge', '0%', '10%', 'Industry-specific expertise'],
    ]
    t2 = Table(score_data, colWidths=[3.5*cm, 3*cm, 3.5*cm, 5*cm])
    t2.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), HexColor('#0B1F3B')),
        ('TEXTCOLOR', (0, 0), (-1, 0), HexColor('#FFFFFF')),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, -1), 9),
        ('GRID', (0, 0), (-1, -1), 0.5, HexColor('#CCCCCC')),
        ('ROWBACKGROUNDS', (0, 1), (-1, -1), [HexColor('#FFFFFF'), HexColor('#F0F5FC')]),
        ('TOPPADDING', (0, 0), (-1, -1), 4),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
    ]))
    story.append(t2)

    # Milestones
    story.append(Spacer(1, 1*cm))
    story.append(Paragraph("4. Project Milestones", styles['CustomH1']))
    mile_data = [
        ['Milestone', 'Date', 'Confidence'],
        ['Review Demo Ready', 'Mar 1, 2026', '85%'],
        ['MVP Complete', 'Mar 25, 2026', '75%'],
        ['Feature Complete', 'Apr 15, 2026', '70%'],
        ['Production Stable', 'May 6, 2026', '65%'],
    ]
    t3 = Table(mile_data, colWidths=[5*cm, 5*cm, 5*cm])
    t3.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), HexColor('#0B1F3B')),
        ('TEXTCOLOR', (0, 0), (-1, 0), HexColor('#FFFFFF')),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, -1), 10),
        ('GRID', (0, 0), (-1, -1), 0.5, HexColor('#CCCCCC')),
        ('ROWBACKGROUNDS', (0, 1), (-1, -1), [HexColor('#FFFFFF'), HexColor('#F0F5FC')]),
        ('TOPPADDING', (0, 0), (-1, -1), 5),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 5),
    ]))
    story.append(t3)

    try:
        doc.build(story)
        print(f"  PDF saved: {out_path} ({os.path.getsize(out_path) // 1024} KB)")
    except Exception as e:
        print(f"  [ERROR] PDF generation failed: {e}")


# ═══════════════════════════════════════════════════════════════
#  MAIN ENTRY POINT
# ═══════════════════════════════════════════════════════════════
if __name__ == "__main__":
    print("=" * 65)
    print("  Smart Hiring — Review Document Generator")
    print("  Output: DOCX + PPTX + PDF")
    print("=" * 65)

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    DIAG_DIR.mkdir(parents=True, exist_ok=True)

    # Step 1: Render Mermaid diagrams
    print("\n[1/4] Rendering Mermaid diagrams...")
    diag_paths = render_mermaid_diagrams()
    print(f"  {len(diag_paths)} diagrams processed")

    # Step 2: Generate DOCX
    print("\n[2/4] Generating DOCX...")
    generate_review_docx(diag_paths, OUTPUT_DIR / "Smart_Hiring_Review.docx")

    # Step 3: Generate PPTX
    print("\n[3/4] Generating PPTX...")
    generate_review_pptx(diag_paths, OUTPUT_DIR / "Smart_Hiring_Review.pptx")

    # Step 4: Generate PDF
    print("\n[4/4] Generating PDF...")
    generate_review_pdf(diag_paths, OUTPUT_DIR / "Smart_Hiring_Review.pdf")

    print("\n" + "=" * 65)
    print("  ALL REVIEW DOCUMENTS GENERATED")
    print(f"  Output directory: {OUTPUT_DIR}")
    print("=" * 65)
    for f in sorted(OUTPUT_DIR.glob("Smart_Hiring_Review.*")):
        if f.is_file():
            print(f"    {f.name}  ({os.path.getsize(f) // 1024} KB)")
