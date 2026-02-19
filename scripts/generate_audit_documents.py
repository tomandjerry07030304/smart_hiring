#!/usr/bin/env python3
"""
Smart Hiring Enterprise Audit Report — Professional Document Generator
======================================================================
Generates: PDF, DOCX, PPTX with watermarks, rendered diagrams, color-coded tables.
"""

import os, sys, re, base64, textwrap, zlib, struct, math
from pathlib import Path
from datetime import datetime
from io import BytesIO
import urllib.request, urllib.error

# ── Third-party imports ──────────────────────────────────────────────
from fpdf import FPDF
from docx import Document
from docx.shared import Inches, Pt, Cm, RGBColor, Emu
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn, nsdecls
from docx.oxml import parse_xml
from pptx import Presentation
from pptx.util import Inches as PI, Pt as PP, Emu as PE
from pptx.dml.color import RGBColor as PC
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image, ImageDraw, ImageFont

# ── Configuration ────────────────────────────────────────────────────
BASE_DIR  = Path(r"c:\Users\venkat anand\OneDrive\Desktop\@smart\smart_hiring")
OUTPUT_DIR = BASE_DIR / "audit_output"
DIAG_DIR  = OUTPUT_DIR / "diagrams"
MD_FILE   = BASE_DIR / "SMART_HIRING_FULL_AUDIT.md"

NAVY       = (0, 32, 63)
DARK_BLUE  = (0, 51, 102)
STEEL      = (70, 130, 180)
LIGHT_BLUE = (200, 220, 240)
WHITE      = (255, 255, 255)
RED        = (220, 53, 69)
ORANGE     = (255, 152, 0)
GREEN      = (40, 167, 69)
YELLOW     = (255, 193, 7)
LGRAY      = (245, 245, 245)
DGRAY      = (60, 60, 60)

# ── Mermaid diagram code blocks ──────────────────────────────────────
MERMAID_BLOCKS = {}  # populated by parser

# ═══════════════════════════════════════════════════════════════════════
#  MARKDOWN PARSER
# ═══════════════════════════════════════════════════════════════════════
def parse_markdown(filepath):
    """Parse markdown into structured list of elements."""
    elements = []
    para_buf = []
    table_buf = []
    mermaid_buf = []
    code_buf = []
    in_mermaid = False
    in_code = False
    diag_idx = 0

    def flush_para():
        nonlocal para_buf
        if para_buf:
            elements.append({"type": "paragraph", "text": " ".join(para_buf)})
            para_buf = []

    def flush_table():
        nonlocal table_buf
        if not table_buf:
            return
        headers = table_buf[0]
        rows = table_buf[1:]
        elements.append({"type": "table", "headers": headers, "rows": rows})
        table_buf = []

    with open(filepath, "r", encoding="utf-8") as f:
        lines = f.readlines()

    i = 0
    while i < len(lines):
        raw = lines[i].rstrip("\n")
        stripped = raw.strip()

        # ── Mermaid fence ────────────────────────────
        if stripped == "```mermaid":
            flush_para(); flush_table()
            in_mermaid = True; mermaid_buf = []; i += 1; continue
        if in_mermaid:
            if stripped == "```":
                in_mermaid = False
                code = "\n".join(mermaid_buf)
                diag_idx += 1
                key = f"diagram_{diag_idx}"
                MERMAID_BLOCKS[key] = code
                elements.append({"type": "mermaid", "key": key})
            else:
                mermaid_buf.append(raw)
            i += 1; continue

        # ── Generic code fence ───────────────────────
        if stripped.startswith("```") and not in_code:
            flush_para(); flush_table()
            in_code = True; code_buf = []; i += 1; continue
        if in_code:
            if stripped == "```":
                in_code = False
                elements.append({"type": "code", "text": "\n".join(code_buf)})
            else:
                code_buf.append(raw)
            i += 1; continue

        # ── Heading ──────────────────────────────────
        m = re.match(r"^(#{1,4})\s+(.*)", raw)
        if m:
            flush_para(); flush_table()
            elements.append({"type": "heading", "level": len(m.group(1)), "text": m.group(2).strip()})
            i += 1; continue

        # ── Horizontal rule ──────────────────────────
        if stripped == "---":
            flush_para(); flush_table()
            elements.append({"type": "hr"})
            i += 1; continue

        # ── Table row ────────────────────────────────
        if stripped.startswith("|") and stripped.endswith("|"):
            flush_para()
            cells = [c.strip() for c in stripped.split("|")[1:-1]]
            # skip separator rows
            if all(re.match(r"^[-:]+$", c) for c in cells):
                i += 1; continue
            table_buf.append(cells)
            i += 1; continue
        else:
            flush_table()

        # ── List item ────────────────────────────────
        lm = re.match(r"^\s*[-*]\s+(.*)", raw)
        nm = re.match(r"^\s*\d+\.\s+(.*)", raw)
        if lm:
            flush_para()
            elements.append({"type": "list_item", "text": lm.group(1)})
            i += 1; continue
        if nm:
            flush_para()
            elements.append({"type": "list_item", "text": nm.group(1)})
            i += 1; continue

        # ── Blank line ───────────────────────────────
        if not stripped:
            flush_para()
            i += 1; continue

        # ── Paragraph text ───────────────────────────
        para_buf.append(raw.strip())
        i += 1

    flush_para(); flush_table()
    return elements


def clean(t):
    """Strip markdown bold/italic/backtick markers for plain-text rendering."""
    t = re.sub(r"\*\*(.+?)\*\*", r"\1", t)
    t = re.sub(r"\*(.+?)\*", r"\1", t)
    t = re.sub(r"`(.+?)`", r"\1", t)
    t = t.replace("\u2705", "[OK]").replace("\u26a0\ufe0f", "[!]").replace("\u274c", "[X]")
    t = t.replace("\u26a0", "[!]")
    t = t.replace("\u2b50", "*").replace("\u2192", "->").replace("\u2265", ">=")
    t = t.replace("\u2014", " - ").replace("\u2013", "-")
    t = t.replace("\u2026", "...").replace("\u2500", "-")
    # Remove any remaining non-latin1 characters for Helvetica compatibility
    try:
        t.encode('latin-1')
    except UnicodeEncodeError:
        t = t.encode('latin-1', errors='replace').decode('latin-1')
    return t


# ═══════════════════════════════════════════════════════════════════════
#  MERMAID RENDERER
# ═══════════════════════════════════════════════════════════════════════
def render_all_diagrams():
    """Render every Mermaid block to PNG via kroki.io API."""
    DIAG_DIR.mkdir(parents=True, exist_ok=True)
    paths = {}
    for key, code in MERMAID_BLOCKS.items():
        out = DIAG_DIR / f"{key}.png"
        if out.exists():
            paths[key] = str(out)
            print(f"  [cached] {key}")
            continue
        ok = _render_one(code, out, key)
        paths[key] = str(out)
    return paths


def _render_one(code, out_path, title):
    """Attempt mermaid.ink, fall back to placeholder."""
    try:
        # Use mermaid.ink API
        encoded = base64.urlsafe_b64encode(code.encode("utf-8")).decode("ascii")
        url = f"https://mermaid.ink/img/{encoded}?type=png&bgColor=white"
        req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
        resp = urllib.request.urlopen(req, timeout=90)
        data = resp.read()
        with open(out_path, "wb") as f:
            f.write(data)
        img = Image.open(out_path)
        print(f"  [rendered] {title}  ({img.width}x{img.height})")
        return True
    except Exception as e:
        print(f"  [fallback] {title}: {e}")
        _placeholder(title, out_path)
        return False


def _placeholder(title, out_path, w=900, h=350):
    img = Image.new("RGB", (w, h), (245, 247, 250))
    d = ImageDraw.Draw(img)
    d.rectangle([3, 3, w - 4, h - 4], outline=(180, 190, 200), width=2)
    try:
        fnt = ImageFont.truetype("arial.ttf", 18)
        sfnt = ImageFont.truetype("arial.ttf", 13)
    except Exception:
        fnt = ImageFont.load_default()
        sfnt = fnt
    d.text((w // 2, h // 2 - 15), "[Architecture Diagram]", fill=(80, 80, 80), font=fnt, anchor="mm")
    d.text((w // 2, h // 2 + 20), title, fill=(130, 130, 130), font=sfnt, anchor="mm")
    img.save(out_path, "PNG")


# ═══════════════════════════════════════════════════════════════════════
#  PDF GENERATOR  (fpdf2)
# ═══════════════════════════════════════════════════════════════════════
class AuditPDF(FPDF):
    """Custom PDF with watermark, professional headers/footers."""

    def __init__(self):
        super().__init__(orientation="P", unit="mm", format="A4")
        self.set_auto_page_break(auto=True, margin=28)
        self.title_page_done = False
        self._toc_entries = []
        self._link_page = {}
        # Try to add Arial/DejaVu for Unicode fallback
        try:
            self.add_font("DejaVu", "", os.path.join(os.path.dirname(__file__), "fonts", "DejaVuSans.ttf"))
        except Exception:
            pass

    # ── header / footer ──────────────────────────────
    def header(self):
        if self.page_no() <= 2:
            return
        # watermark
        self.set_font("Helvetica", "B", 38)
        self.set_text_color(235, 235, 235)
        with self.rotation(45, self.w / 2, self.h / 2):
            self.set_xy(self.w / 2 - 65, self.h / 2)
            self.cell(130, 12, "CONFIDENTIAL", align="C")
        # top bar
        self.set_fill_color(*NAVY)
        self.rect(0, 0, self.w, 12, "F")
        self.set_font("Helvetica", "B", 7)
        self.set_text_color(*WHITE)
        self.set_xy(10, 3)
        self.cell(0, 6, "CONFIDENTIAL  |  INTERNAL AUDIT  |  Smart Hiring Enterprise Audit Report", align="C")
        self.set_text_color(*DGRAY)

    def footer(self):
        if self.page_no() <= 1:
            return
        self.set_y(-18)
        self.set_draw_color(*NAVY)
        self.line(15, self.h - 18, self.w - 15, self.h - 18)
        self.set_font("Helvetica", "I", 7)
        self.set_text_color(120, 120, 120)
        self.set_y(-14)
        self.cell(0, 8, f"Smart Hiring Enterprise Audit Report  |  v1.0  |  February 2026  |  Page {self.page_no()}", align="C")

    # ── title page ───────────────────────────────────
    def make_title_page(self):
        self.add_page()
        # dark header block
        self.set_fill_color(*NAVY)
        self.rect(0, 0, self.w, 105, "F")
        # accent line
        self.set_fill_color(*STEEL)
        self.rect(0, 105, self.w, 3, "F")

        self.set_text_color(*WHITE)
        self.set_font("Helvetica", "B", 28)
        self.set_xy(20, 25)
        self.multi_cell(self.w - 40, 12, "Smart Hiring\nEnterprise Audit Report", align="L")
        self.set_font("Helvetica", "", 13)
        self.set_xy(20, 65)
        self.cell(0, 8, "System: Smart Hiring System  -  Enterprise Edition v4.2+")
        self.set_xy(20, 75)
        self.cell(0, 8, "Audit Classification: CONFIDENTIAL")
        self.set_xy(20, 85)
        self.cell(0, 8, "Date: February 12, 2026   |   Version: 1.0")

        # body
        self.set_text_color(*DGRAY)
        self.set_font("Helvetica", "", 10)
        y = 120
        info = [
            ("Audit Level:", "CTO / Principal Architect / ML Lead / Security Auditor / Compliance Officer"),
            ("Prepared By:", "Enterprise Architecture & AI Governance Team"),
            ("Distribution:", "CTO, VP Engineering, ML Lead, Security Team"),
            ("Next Review:", "May 2026 (Quarterly)"),
        ]
        for label, val in info:
            self.set_font("Helvetica", "B", 10)
            self.set_xy(25, y)
            self.cell(40, 7, label)
            self.set_font("Helvetica", "", 10)
            self.cell(0, 7, val)
            y += 10

        # signature block
        y = 200
        self.set_font("Helvetica", "B", 11)
        self.set_xy(25, y)
        self.cell(0, 8, "APPROVAL SIGNATURES")
        self.set_draw_color(*NAVY)
        self.line(25, y + 10, 185, y + 10)
        sigs = [
            ("CTO Approval:", "___________________________    Date: ___________"),
            ("Lead Auditor:", "___________________________    Date: ___________"),
            ("Security Lead:", "___________________________    Date: ___________"),
        ]
        y += 18
        for label, line in sigs:
            self.set_font("Helvetica", "B", 9)
            self.set_xy(25, y)
            self.cell(35, 7, label)
            self.set_font("Helvetica", "", 9)
            self.cell(0, 7, line)
            y += 16

        self.title_page_done = True

    # ── table of contents ────────────────────────────
    def make_toc(self, elements):
        self.add_page()
        self.set_font("Helvetica", "B", 20)
        self.set_text_color(*NAVY)
        self.cell(0, 14, "Table of Contents", new_x="LMARGIN", new_y="NEXT")
        self.set_draw_color(*STEEL)
        self.line(15, self.get_y(), 195, self.get_y())
        self.ln(6)
        self.set_text_color(*DGRAY)
        sec = 0
        for el in elements:
            if el["type"] == "heading" and el["level"] <= 2:
                lvl = el["level"]
                txt = clean(el["text"])
                if lvl == 1:
                    self.set_font("Helvetica", "B", 11)
                    indent = 0
                else:
                    self.set_font("Helvetica", "", 10)
                    indent = 8
                self.set_x(15 + indent)
                self.cell(0, 7, txt, new_x="LMARGIN", new_y="NEXT")
                if self.get_y() > 270:
                    self.add_page()

    # ── render elements ──────────────────────────────
    def render_elements(self, elements, diag_paths):
        for el in elements:
            tp = el["type"]
            if tp == "heading":
                self._render_heading(el)
            elif tp == "paragraph":
                self._render_para(el)
            elif tp == "table":
                self._render_table(el)
            elif tp == "list_item":
                self._render_list(el)
            elif tp == "mermaid":
                self._render_diagram(el, diag_paths)
            elif tp == "code":
                self._render_code(el)
            elif tp == "hr":
                self._render_hr()

    def _render_heading(self, el):
        lvl = el["level"]
        txt = clean(el["text"])
        if lvl == 1:
            self.add_page()
            # section header bar
            self.set_fill_color(*NAVY)
            self.rect(15, self.get_y(), self.w - 30, 12, "F")
            self.set_font("Helvetica", "B", 16)
            self.set_text_color(*WHITE)
            self.set_xy(20, self.get_y() + 1)
            self.cell(self.w - 40, 10, txt, align="L")
            self.set_text_color(*DGRAY)
            self.ln(16)
        elif lvl == 2:
            if self.get_y() > 240:
                self.add_page()
            self.ln(4)
            self.set_font("Helvetica", "B", 13)
            self.set_text_color(*DARK_BLUE)
            self.cell(0, 9, txt, new_x="LMARGIN", new_y="NEXT")
            self.set_draw_color(*STEEL)
            self.line(15, self.get_y(), 120, self.get_y())
            self.set_text_color(*DGRAY)
            self.ln(3)
        elif lvl == 3:
            if self.get_y() > 250:
                self.add_page()
            self.ln(2)
            self.set_font("Helvetica", "B", 11)
            self.set_text_color(*NAVY)
            self.cell(0, 8, txt, new_x="LMARGIN", new_y="NEXT")
            self.set_text_color(*DGRAY)
            self.ln(1)
        else:
            self.set_font("Helvetica", "BI", 10)
            self.set_text_color(*STEEL)
            self.cell(0, 7, txt, new_x="LMARGIN", new_y="NEXT")
            self.set_text_color(*DGRAY)
            self.ln(1)

    def _render_para(self, el):
        txt = clean(el["text"])
        if not txt.strip():
            return
        self.set_font("Helvetica", "", 9)
        self.set_text_color(*DGRAY)
        self.set_x(15)
        # Handle bold segments within paragraph
        self.multi_cell(self.w - 30, 5, txt, new_x="LMARGIN", new_y="NEXT")
        self.ln(2)

    def _render_table(self, el):
        headers = el["headers"]
        rows = el["rows"]
        if not headers:
            return
        ncols = len(headers)
        avail = self.w - 30
        # calculate column widths
        col_ws = [avail / ncols] * ncols
        # Try smarter sizing: measure text
        for ci, h in enumerate(headers):
            max_len = len(clean(h))
            for r in rows:
                if ci < len(r):
                    max_len = max(max_len, len(clean(r[ci])))
            col_ws[ci] = max_len
        total = sum(col_ws)
        col_ws = [max(w / total * avail, 18) for w in col_ws]
        # re-normalize
        s = sum(col_ws)
        col_ws = [w / s * avail for w in col_ws]

        if self.get_y() + 12 + len(rows) * 7 > 270:
            self.add_page()

        # header row
        self.set_fill_color(*NAVY)
        self.set_text_color(*WHITE)
        self.set_font("Helvetica", "B", 7.5)
        x0 = 15
        y0 = self.get_y()
        max_h = 6
        for ci, h in enumerate(headers):
            self.set_xy(x0, y0)
            self.cell(col_ws[ci], max_h, clean(h)[:50], border=1, fill=True, align="C")
            x0 += col_ws[ci]
        self.ln(max_h)

        # data rows
        self.set_text_color(*DGRAY)
        self.set_font("Helvetica", "", 7)
        for ri, row in enumerate(rows):
            if self.get_y() > 268:
                self.add_page()
                # re-print header
                self.set_fill_color(*NAVY)
                self.set_text_color(*WHITE)
                self.set_font("Helvetica", "B", 7.5)
                x0 = 15
                y0 = self.get_y()
                for ci, h in enumerate(headers):
                    self.set_xy(x0, y0)
                    self.cell(col_ws[ci], 6, clean(h)[:50], border=1, fill=True, align="C")
                    x0 += col_ws[ci]
                self.ln(6)
                self.set_text_color(*DGRAY)
                self.set_font("Helvetica", "", 7)

            fill = ri % 2 == 0
            if fill:
                self.set_fill_color(*LGRAY)
            x0 = 15
            y0 = self.get_y()
            for ci in range(ncols):
                val = clean(row[ci]) if ci < len(row) else ""
                # color-code severity cells
                cell_fill = fill
                orig_fill = LGRAY
                sev = val.upper()
                if "CRITICAL" in sev or "P0" in sev:
                    self.set_fill_color(255, 200, 200)
                    cell_fill = True
                elif "HIGH" in sev or "P1" in sev:
                    self.set_fill_color(255, 230, 200)
                    cell_fill = True
                elif "MEDIUM" in sev:
                    self.set_fill_color(255, 245, 200)
                    cell_fill = True
                elif val.startswith("[OK]"):
                    self.set_fill_color(200, 240, 200)
                    cell_fill = True
                elif val.startswith("[X]"):
                    self.set_fill_color(255, 200, 200)
                    cell_fill = True
                elif val.startswith("[!]"):
                    self.set_fill_color(255, 230, 200)
                    cell_fill = True

                self.set_xy(x0, y0)
                self.cell(col_ws[ci], 6, val[:80], border=1, fill=cell_fill, align="L")
                x0 += col_ws[ci]
                if fill:
                    self.set_fill_color(*LGRAY)
                else:
                    self.set_fill_color(*WHITE)
            self.ln(6)
        self.ln(3)

    def _render_list(self, el):
        txt = clean(el["text"])
        self.set_font("Helvetica", "", 9)
        self.set_text_color(*DGRAY)
        self.set_x(20)
        self.cell(5, 5, "-")
        self.multi_cell(self.w - 40, 5, txt, new_x="LMARGIN", new_y="NEXT")
        self.ln(1)

    def _render_diagram(self, el, diag_paths):
        key = el["key"]
        path = diag_paths.get(key)
        if not path or not os.path.exists(path):
            return
        self.add_page()
        try:
            img = Image.open(path)
            iw, ih = img.size
            max_w = self.w - 30
            max_h = self.h - 60
            ratio = min(max_w / (iw * 0.264583), max_h / (ih * 0.264583))
            w_mm = iw * 0.264583 * ratio
            h_mm = ih * 0.264583 * ratio
            x = (self.w - w_mm) / 2
            self.image(path, x=x, y=self.get_y(), w=w_mm)
            self.set_y(self.get_y() + h_mm + 5)
        except Exception as e:
            self.set_font("Helvetica", "I", 9)
            self.cell(0, 8, f"[Diagram: {key}]", new_x="LMARGIN", new_y="NEXT")

    def _render_code(self, el):
        txt = clean(el["text"])
        if not txt.strip():
            return
        self.set_fill_color(245, 245, 248)
        self.set_font("Courier", "", 7)
        self.set_text_color(50, 50, 50)
        lines = txt.split("\n")
        y0 = self.get_y()
        block_h = len(lines) * 4 + 4
        if y0 + block_h > 270:
            self.add_page()
        self.set_x(15)
        self.rect(15, self.get_y(), self.w - 30, block_h, "F")
        self.ln(2)
        for line in lines:
            self.set_x(18)
            self.cell(0, 4, line[:120], new_x="LMARGIN", new_y="NEXT")
        self.ln(3)
        self.set_text_color(*DGRAY)

    def _render_hr(self):
        self.ln(3)
        self.set_draw_color(200, 200, 200)
        self.line(15, self.get_y(), self.w - 15, self.get_y())
        self.ln(3)


def generate_pdf(elements, diag_paths, out_path):
    print("  Creating PDF...")
    pdf = AuditPDF()
    pdf.set_title("Smart Hiring Enterprise Audit Report")
    pdf.set_author("Enterprise Architecture & AI Governance Team")
    pdf.make_title_page()
    pdf.make_toc(elements)
    pdf.render_elements(elements, diag_paths)
    pdf.output(str(out_path))
    print(f"  PDF saved: {out_path}")
    sz = os.path.getsize(out_path)
    print(f"  Size: {sz / 1024:.0f} KB  |  Pages: {pdf.page_no()}")


# ═══════════════════════════════════════════════════════════════════════
#  DOCX GENERATOR  (python-docx)
# ═══════════════════════════════════════════════════════════════════════
def generate_docx(elements, diag_paths, out_path):
    print("  Creating DOCX...")
    doc = Document()

    # ── Styles ───────────────────────────────────────
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(10)
    style.font.color.rgb = RGBColor(50, 50, 50)
    style.paragraph_format.space_after = Pt(4)

    for lvl in range(1, 5):
        hs = doc.styles[f"Heading {lvl}"]
        hs.font.name = "Calibri"
        hs.font.color.rgb = RGBColor(*NAVY)
        if lvl == 1:
            hs.font.size = Pt(20)
        elif lvl == 2:
            hs.font.size = Pt(15)
        elif lvl == 3:
            hs.font.size = Pt(12)
        else:
            hs.font.size = Pt(11)

    # ── Header with watermark text ───────────────────
    section = doc.sections[0]
    section.top_margin = Cm(2)
    section.bottom_margin = Cm(2)
    section.left_margin = Cm(2.5)
    section.right_margin = Cm(2.5)
    header = section.header
    hp = header.paragraphs[0]
    hp.text = "CONFIDENTIAL — INTERNAL AUDIT  |  Smart Hiring Enterprise Audit Report  |  v1.0"
    hp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    hr = hp.runs[0] if hp.runs else hp.add_run()
    hr.font.size = Pt(7)
    hr.font.color.rgb = RGBColor(150, 150, 150)

    footer = section.footer
    fp = footer.paragraphs[0]
    fp.text = "CONFIDENTIAL  |  Distribution: CTO, VP Engineering, ML Lead, Security Team"
    fp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    fr = fp.runs[0] if fp.runs else fp.add_run()
    fr.font.size = Pt(7)
    fr.font.color.rgb = RGBColor(150, 150, 150)

    # ── Title page ───────────────────────────────────
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.space_before = Pt(80)
    r = p.add_run("SMART HIRING\nENTERPRISE AUDIT REPORT")
    r.bold = True
    r.font.size = Pt(28)
    r.font.color.rgb = RGBColor(*NAVY)

    p2 = doc.add_paragraph()
    p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r2 = p2.add_run("System: Smart Hiring System — Enterprise Edition v4.2+\n"
                     "Audit Classification: CONFIDENTIAL\n"
                     "Date: February 12, 2026  |  Version: 1.0")
    r2.font.size = Pt(11)
    r2.font.color.rgb = RGBColor(80, 80, 80)

    p3 = doc.add_paragraph()
    p3.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p3.space_before = Pt(30)
    r3 = p3.add_run("Prepared By: Enterprise Architecture & AI Governance Team")
    r3.font.size = Pt(10)
    r3.font.color.rgb = RGBColor(100, 100, 100)

    # Signature block
    doc.add_paragraph()
    sig_table = doc.add_table(rows=4, cols=3)
    sig_table.style = "Table Grid"
    sig_hdr = sig_table.rows[0]
    for ci, txt in enumerate(["Role", "Signature", "Date"]):
        cell = sig_hdr.cells[ci]
        cell.text = txt
        _shade_cell(cell, NAVY)
        for run in cell.paragraphs[0].runs:
            run.font.color.rgb = RGBColor(*WHITE)
            run.bold = True
            run.font.size = Pt(9)
    for ri, role in enumerate(["CTO Approval", "Lead Auditor", "Security Lead"], 1):
        sig_table.rows[ri].cells[0].text = role
        sig_table.rows[ri].cells[0].paragraphs[0].runs[0].font.size = Pt(9)

    doc.add_page_break()

    # ── TOC placeholder ──────────────────────────────
    p_toc = doc.add_paragraph()
    p_toc.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r_toc = p_toc.add_run("TABLE OF CONTENTS")
    r_toc.bold = True
    r_toc.font.size = Pt(18)
    r_toc.font.color.rgb = RGBColor(*NAVY)
    doc.add_paragraph()
    # Add TOC field
    p_field = doc.add_paragraph()
    _add_toc_field(p_field)
    doc.add_page_break()

    # ── Render elements ──────────────────────────────
    for el in elements:
        tp = el["type"]
        if tp == "heading":
            lvl = min(el["level"], 4)
            doc.add_heading(clean(el["text"]), level=lvl)
        elif tp == "paragraph":
            txt = clean(el["text"])
            if txt.strip():
                doc.add_paragraph(txt)
        elif tp == "table":
            _docx_table(doc, el)
        elif tp == "list_item":
            doc.add_paragraph(clean(el["text"]), style="List Bullet")
        elif tp == "mermaid":
            _docx_diagram(doc, el, diag_paths)
        elif tp == "code":
            _docx_code(doc, el)
        elif tp == "hr":
            doc.add_paragraph("─" * 80)

    doc.save(str(out_path))
    print(f"  DOCX saved: {out_path}")
    sz = os.path.getsize(out_path)
    print(f"  Size: {sz / 1024:.0f} KB")


def _shade_cell(cell, color):
    """Apply background shading to a DOCX table cell."""
    shading = parse_xml(
        f'<w:shd {nsdecls("w")} w:fill="{_rgb_hex(color)}" w:val="clear"/>'
    )
    cell._tc.get_or_add_tcPr().append(shading)


def _rgb_hex(color):
    return f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"


def _add_toc_field(paragraph):
    """Insert TOC field code into the DOCX."""
    run = paragraph.add_run()
    fld_char_begin = parse_xml(f'<w:fldChar {nsdecls("w")} w:fldCharType="begin"/>')
    run._r.append(fld_char_begin)
    run2 = paragraph.add_run()
    instr = parse_xml(f'<w:instrText {nsdecls("w")} xml:space="preserve"> TOC \\o "1-3" \\h \\z \\u </w:instrText>')
    run2._r.append(instr)
    run3 = paragraph.add_run()
    fld_char_end = parse_xml(f'<w:fldChar {nsdecls("w")} w:fldCharType="end"/>')
    run3._r.append(fld_char_end)
    paragraph.add_run("\n[Right-click and select 'Update Field' to generate TOC]").font.color.rgb = RGBColor(150, 150, 150)


def _docx_table(doc, el):
    headers = el["headers"]
    rows = el["rows"]
    ncols = len(headers)
    t = doc.add_table(rows=1 + len(rows), cols=ncols)
    t.style = "Table Grid"
    t.alignment = WD_TABLE_ALIGNMENT.CENTER
    # header row
    for ci, h in enumerate(headers):
        cell = t.rows[0].cells[ci]
        cell.text = clean(h)
        _shade_cell(cell, NAVY)
        for run in cell.paragraphs[0].runs:
            run.font.color.rgb = RGBColor(*WHITE)
            run.bold = True
            run.font.size = Pt(8)
        cell.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
    # data rows
    for ri, row in enumerate(rows):
        for ci in range(ncols):
            cell = t.rows[ri + 1].cells[ci]
            val = clean(row[ci]) if ci < len(row) else ""
            cell.text = val
            for run in cell.paragraphs[0].runs:
                run.font.size = Pt(8)
            # color code
            up = val.upper()
            if "CRITICAL" in up or "P0" in up:
                _shade_cell(cell, (255, 210, 210))
            elif "HIGH" in up or "P1" in up:
                _shade_cell(cell, (255, 235, 210))
            elif "MEDIUM" in up:
                _shade_cell(cell, (255, 248, 210))
            elif val.startswith("[OK]"):
                _shade_cell(cell, (210, 245, 210))
            elif val.startswith("[X]"):
                _shade_cell(cell, (255, 210, 210))
            elif val.startswith("[!]"):
                _shade_cell(cell, (255, 235, 210))
            elif ri % 2 == 0:
                _shade_cell(cell, (245, 248, 252))
    doc.add_paragraph()


def _docx_diagram(doc, el, diag_paths):
    key = el["key"]
    path = diag_paths.get(key)
    if path and os.path.exists(path):
        try:
            img = Image.open(path)
            iw, ih = img.size
            max_w_in = 6.0
            ratio = max_w_in / (iw / 96)
            w_in = min(iw / 96 * ratio, max_w_in)
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            r = p.add_run()
            r.add_picture(path, width=Inches(w_in))
            doc.add_paragraph()
        except Exception:
            doc.add_paragraph(f"[Diagram: {key}]")
    else:
        doc.add_paragraph(f"[Diagram: {key}]")


def _docx_code(doc, el):
    txt = clean(el["text"])
    if not txt.strip():
        return
    p = doc.add_paragraph()
    r = p.add_run(txt)
    r.font.name = "Consolas"
    r.font.size = Pt(8)
    r.font.color.rgb = RGBColor(50, 50, 50)


# ═══════════════════════════════════════════════════════════════════════
#  PPTX GENERATOR  (python-pptx)
# ═══════════════════════════════════════════════════════════════════════
def generate_pptx(elements, diag_paths, out_path):
    print("  Creating PPTX...")
    prs = Presentation()
    prs.slide_width = PI(13.333)
    prs.slide_height = PI(7.5)

    # ── Helper ───────────────────────────────────────
    def add_slide():
        layout = prs.slide_layouts[6]  # blank
        return prs.slides.add_slide(layout)

    def add_textbox(slide, left, top, width, height):
        return slide.shapes.add_textbox(PI(left), PI(top), PI(width), PI(height))

    def add_bg(slide, color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = PC(*color)

    def add_rect(slide, l, t, w, h, color):
        shape = slide.shapes.add_shape(1, PI(l), PI(t), PI(w), PI(h))  # MSO_SHAPE.RECTANGLE
        shape.fill.solid()
        shape.fill.fore_color.rgb = PC(*color)
        shape.line.fill.background()
        return shape

    # ── 1. Title Slide ───────────────────────────────
    s = add_slide()
    add_bg(s, NAVY)
    # Accent stripe
    add_rect(s, 0, 3.5, 13.333, 0.06, STEEL)

    tb = add_textbox(s, 1, 1.0, 11, 2.0)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Smart Hiring"
    p.font.size = PP(44)
    p.font.bold = True
    p.font.color.rgb = PC(*WHITE)
    p = tf.add_paragraph()
    p.text = "Enterprise Audit Report"
    p.font.size = PP(32)
    p.font.color.rgb = PC(*STEEL)

    tb2 = add_textbox(s, 1, 4.0, 11, 2.0)
    tf2 = tb2.text_frame; tf2.word_wrap = True
    for txt in ["System: Smart Hiring System — Enterprise Edition v4.2+",
                "Audit Classification: CONFIDENTIAL",
                "Date: February 12, 2026  |  Version: 1.0",
                "Prepared By: Enterprise Architecture & AI Governance Team"]:
        p = tf2.add_paragraph()
        p.text = txt
        p.font.size = PP(14)
        p.font.color.rgb = PC(180, 200, 220)
        p.space_after = PP(4)

    # ── 2. Agenda Slide ──────────────────────────────
    s = add_slide()
    add_bg(s, WHITE)
    add_rect(s, 0, 0, 13.333, 0.8, NAVY)
    tb = add_textbox(s, 0.5, 0.1, 12, 0.6)
    p = tb.text_frame.paragraphs[0]
    p.text = "AGENDA"
    p.font.size = PP(28); p.font.bold = True; p.font.color.rgb = PC(*WHITE)

    agenda = [
        "1. Executive Summary & Critical Findings",
        "2. System Architecture Overview",
        "3. ML Model Inventory & Pipeline",
        "4. Ranking & Fairness Engine",
        "5. Security & Compliance Audit",
        "6. Performance & Scalability",
        "7. Risk Assessment Matrix",
        "8. Enterprise Roadmap & Recommendations",
    ]
    tb2 = add_textbox(s, 1.5, 1.3, 10, 5.5)
    tf2 = tb2.text_frame; tf2.word_wrap = True
    for item in agenda:
        p = tf2.add_paragraph()
        p.text = item
        p.font.size = PP(18)
        p.font.color.rgb = PC(*DGRAY)
        p.space_after = PP(10)

    # ── 3. Overall Assessment Slide ──────────────────
    s = add_slide()
    add_bg(s, WHITE)
    add_rect(s, 0, 0, 13.333, 0.8, NAVY)
    tb = add_textbox(s, 0.5, 0.1, 12, 0.6)
    p = tb.text_frame.paragraphs[0]
    p.text = "OVERALL ASSESSMENT"; p.font.size = PP(28); p.font.bold = True; p.font.color.rgb = PC(*WHITE)

    scores = [
        ("Architecture Maturity", "7.2/10", "Production-Ready", GREEN),
        ("ML Pipeline Robustness", "6.8/10", "Functional, needs hardening", YELLOW),
        ("Security Posture", "5.5/10", "Significant gaps", RED),
        ("Fairness & Compliance", "8.0/10", "Strong foundation", GREEN),
        ("Code Quality", "6.0/10", "Moderate", YELLOW),
        ("Scalability", "5.0/10", "Single-node only", ORANGE),
        ("MLOps Maturity", "4.5/10", "Early stage", RED),
    ]
    y = 1.2
    for name, score, rating, color in scores:
        # Background bar
        add_rect(s, 1, y, 11.333, 0.65, (245, 248, 252))
        # Color indicator
        add_rect(s, 1, y, 0.15, 0.65, color)
        # Name
        tb = add_textbox(s, 1.4, y + 0.05, 4.5, 0.55)
        p = tb.text_frame.paragraphs[0]; p.text = name; p.font.size = PP(14); p.font.color.rgb = PC(*DGRAY)
        # Score
        tb = add_textbox(s, 6, y + 0.05, 1.5, 0.55)
        p = tb.text_frame.paragraphs[0]; p.text = score; p.font.size = PP(16); p.font.bold = True
        p.font.color.rgb = PC(*color)
        # Rating
        tb = add_textbox(s, 7.8, y + 0.05, 4.5, 0.55)
        p = tb.text_frame.paragraphs[0]; p.text = rating; p.font.size = PP(12); p.font.color.rgb = PC(120, 120, 120)
        y += 0.78

    # ── 4. Critical Findings Slide ───────────────────
    s = add_slide()
    add_bg(s, WHITE)
    add_rect(s, 0, 0, 13.333, 0.8, RED)
    tb = add_textbox(s, 0.5, 0.1, 12, 0.6)
    p = tb.text_frame.paragraphs[0]
    p.text = "CRITICAL FINDINGS"; p.font.size = PP(28); p.font.bold = True; p.font.color.rgb = PC(*WHITE)

    findings = [
        ("CF-01", "CRITICAL", "JWT secret defaults to hardcoded value", RED),
        ("CF-02", "CRITICAL", "PII encryption module exists but NOT integrated", RED),
        ("CF-03", "HIGH", "Dual queue systems operate independently", ORANGE),
        ("CF-04", "HIGH", "RBAC module built but unused", ORANGE),
        ("CF-05", "HIGH", "Three different scoring weight formulas", ORANGE),
        ("CF-06", "HIGH", "Assessment queries wrong MongoDB collection", ORANGE),
    ]
    y = 1.2
    for fid, sev, desc, color in findings:
        add_rect(s, 1, y, 0.15, 0.72, color)
        # ID badge
        shape = add_rect(s, 1.4, y + 0.08, 1.0, 0.55, color)
        shape.text_frame.paragraphs[0].text = fid
        shape.text_frame.paragraphs[0].font.size = PP(10)
        shape.text_frame.paragraphs[0].font.color.rgb = PC(*WHITE)
        shape.text_frame.paragraphs[0].font.bold = True
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        # Severity
        tb = add_textbox(s, 2.7, y + 0.08, 1.5, 0.55)
        p = tb.text_frame.paragraphs[0]; p.text = sev; p.font.size = PP(11); p.font.bold = True
        p.font.color.rgb = PC(*color)
        # Description
        tb = add_textbox(s, 4.3, y + 0.08, 8.5, 0.55)
        p = tb.text_frame.paragraphs[0]; p.text = desc; p.font.size = PP(13); p.font.color.rgb = PC(*DGRAY)
        y += 0.85

    # remaining findings (medium/low) note
    tb = add_textbox(s, 1, y + 0.3, 11, 0.5)
    p = tb.text_frame.paragraphs[0]
    p.text = "+ 4 additional findings (Medium/Low severity) — see full report"
    p.font.size = PP(11); p.font.italic = True; p.font.color.rgb = PC(130, 130, 130)

    # ── 5. Key Strengths Slide ───────────────────────
    s = add_slide()
    add_bg(s, WHITE)
    add_rect(s, 0, 0, 13.333, 0.8, GREEN)
    tb = add_textbox(s, 0.5, 0.1, 12, 0.6)
    p = tb.text_frame.paragraphs[0]
    p.text = "KEY STRENGTHS"; p.font.size = PP(28); p.font.bold = True; p.font.color.rgb = PC(*WHITE)

    strengths = [
        "Triple-fallback ML pipeline (SBERT -> TF-IDF -> Keyword)",
        "Custom fairness engine with 8 statistical metrics",
        "Context-aware skill extraction with anti-fraud measures",
        "NER-based resume anonymization for bias-free scoring",
        "Explainable AI with GDPR Art. 22 transparency reports",
        "Containerized multi-stage Docker deployment",
        "Comprehensive background task system with DLQ + retry logic",
    ]
    y = 1.2
    for st in strengths:
        add_rect(s, 1, y, 0.12, 0.55, GREEN)
        tb = add_textbox(s, 1.5, y, 11, 0.55)
        p = tb.text_frame.paragraphs[0]; p.text = st; p.font.size = PP(14); p.font.color.rgb = PC(*DGRAY)
        y += 0.72

    # ── 6. Architecture diagram slides ───────────────
    diagram_slides = [
        ("diagram_1", "High-Level System Architecture"),
        ("diagram_2", "Resume Processing Pipeline"),
        ("diagram_3", "ML Ranking Flow"),
        ("diagram_4", "Worker & Queue Architecture"),
        ("diagram_5", "Fairness Engine Pipeline"),
        ("diagram_6", "Deployment Architecture"),
    ]
    for key, title in diagram_slides:
        path = diag_paths.get(key)
        if not path or not os.path.exists(path):
            continue
        s = add_slide()
        add_bg(s, WHITE)
        add_rect(s, 0, 0, 13.333, 0.8, DARK_BLUE)
        tb = add_textbox(s, 0.5, 0.1, 12, 0.6)
        p = tb.text_frame.paragraphs[0]
        p.text = title.upper(); p.font.size = PP(22); p.font.bold = True; p.font.color.rgb = PC(*WHITE)
        try:
            img = Image.open(path)
            iw, ih = img.size
            max_w, max_h = 11.5, 5.8
            ratio = min(max_w / (iw / 96), max_h / (ih / 96))
            w = iw / 96 * ratio
            h = ih / 96 * ratio
            left = (13.333 - w) / 2
            top = 0.9 + (5.8 - h) / 2
            s.shapes.add_picture(path, PI(left), PI(top), PI(w), PI(h))
        except Exception:
            pass

    # ── 7. Technology Stack Slide ────────────────────
    s = add_slide()
    add_bg(s, WHITE)
    add_rect(s, 0, 0, 13.333, 0.8, NAVY)
    tb = add_textbox(s, 0.5, 0.1, 12, 0.6)
    p = tb.text_frame.paragraphs[0]
    p.text = "TECHNOLOGY STACK"; p.font.size = PP(28); p.font.bold = True; p.font.color.rgb = PC(*WHITE)

    stack = [
        ("Backend", "Flask 3.0.0 + Gunicorn + Eventlet"),
        ("Database", "MongoDB 7.0"),
        ("Cache/Queue", "Redis 7.x + Celery 5.3.4"),
        ("ML Embeddings", "SBERT all-MiniLM-L6-v2 (384-dim)"),
        ("NLP", "spaCy en_core_web_sm"),
        ("Fairness", "Custom Engine + Fairlearn"),
        ("Real-Time", "Flask-SocketIO 5.3.5"),
        ("Auth", "Flask-JWT-Extended 4.6.0"),
        ("Deploy", "Docker Multi-stage + docker-compose"),
    ]
    y = 1.2
    left_x = 1.0
    for label, val in stack:
        add_rect(s, left_x, y, 3.0, 0.62, (235, 240, 248))
        tb = add_textbox(s, left_x + 0.2, y + 0.05, 2.6, 0.52)
        p = tb.text_frame.paragraphs[0]; p.text = label; p.font.size = PP(12); p.font.bold = True
        p.font.color.rgb = PC(*NAVY)
        tb = add_textbox(s, left_x + 3.2, y + 0.05, 6, 0.52)
        p = tb.text_frame.paragraphs[0]; p.text = val; p.font.size = PP(12)
        p.font.color.rgb = PC(*DGRAY)
        y += 0.68

    # ── 8. ML Model Inventory Slide ──────────────────
    s = add_slide()
    add_bg(s, WHITE)
    add_rect(s, 0, 0, 13.333, 0.8, DARK_BLUE)
    tb = add_textbox(s, 0.5, 0.1, 12, 0.6)
    p = tb.text_frame.paragraphs[0]
    p.text = "ML MODEL INVENTORY"; p.font.size = PP(28); p.font.bold = True; p.font.color.rgb = PC(*WHITE)

    models = [
        ("all-MiniLM-L6-v2", "Semantic Embedding", "80MB, 384-dim", "Active"),
        ("en_core_web_sm", "NER / Resume Parsing", "12MB", "Active"),
        ("TfidfVectorizer", "TF-IDF Fallback Similarity", "5-20MB", "Active"),
        ("Skill Extractor", "350+ Skill Taxonomy", "2MB", "Active"),
        ("Fairness Engine", "8 Statistical Metrics", "CPU", "Active"),
        ("CCI Calculator", "Career Consistency Index", "CPU", "Active"),
        ("cross-encoder/ms-marco", "Reranker (Planned)", "80MB", "Recommended"),
    ]
    y = 1.1
    # Header row
    add_rect(s, 0.8, y, 11.7, 0.55, NAVY)
    for xi, (col, cx) in enumerate([("Model", 1.0), ("Purpose", 4.2), ("Size", 8.5), ("Status", 10.5)]):
        tb = add_textbox(s, cx, y + 0.03, 3, 0.5)
        p = tb.text_frame.paragraphs[0]; p.text = col; p.font.size = PP(11); p.font.bold = True
        p.font.color.rgb = PC(*WHITE)
    y += 0.6
    for model, purpose, size, status in models:
        bg_color = (245, 248, 252) if models.index((model, purpose, size, status)) % 2 == 0 else WHITE
        add_rect(s, 0.8, y, 11.7, 0.55, bg_color)
        for val, cx in [(model, 1.0), (purpose, 4.2), (size, 8.5), (status, 10.5)]:
            tb = add_textbox(s, cx, y + 0.03, 3.5, 0.5)
            p = tb.text_frame.paragraphs[0]; p.text = val; p.font.size = PP(10)
            if status == "Recommended" and val == status:
                p.font.color.rgb = PC(*STEEL)
                p.font.italic = True
            else:
                p.font.color.rgb = PC(*DGRAY)
        y += 0.58

    # ── 9. Security Audit Slide ──────────────────────
    s = add_slide()
    add_bg(s, WHITE)
    add_rect(s, 0, 0, 13.333, 0.8, RED)
    tb = add_textbox(s, 0.5, 0.1, 12, 0.6)
    p = tb.text_frame.paragraphs[0]
    p.text = "SECURITY & COMPLIANCE SNAPSHOT"; p.font.size = PP(24); p.font.bold = True; p.font.color.rgb = PC(*WHITE)

    controls = [
        ("JWT Auth", "PASS", GREEN),
        ("Password Hashing (bcrypt)", "PASS", GREEN),
        ("PII Encryption at Rest", "FAIL", RED),
        ("RBAC Enforcement", "WARN", ORANGE),
        ("GDPR DSR Routes", "PASS", GREEN),
        ("TLS Certificate Validation", "FAIL", RED),
        ("SSRF Protection", "FAIL", RED),
        ("Rate Limiter (Redis-backed)", "WARN", ORANGE),
        ("Security Headers", "PASS", GREEN),
        ("Fairness Compliance (EEOC)", "PASS", GREEN),
    ]
    y = 1.1
    for name, status, color in controls:
        add_rect(s, 1, y, 0.12, 0.52, color)
        tb = add_textbox(s, 1.4, y, 6, 0.52)
        p = tb.text_frame.paragraphs[0]; p.text = name; p.font.size = PP(13); p.font.color.rgb = PC(*DGRAY)
        shape = add_rect(s, 8.0, y + 0.05, 1.2, 0.42, color)
        shape.text_frame.paragraphs[0].text = status
        shape.text_frame.paragraphs[0].font.size = PP(10)
        shape.text_frame.paragraphs[0].font.bold = True
        shape.text_frame.paragraphs[0].font.color.rgb = PC(*WHITE)
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        y += 0.58

    # ── 10. Risk Matrix Slide ────────────────────────
    s = add_slide()
    add_bg(s, WHITE)
    add_rect(s, 0, 0, 13.333, 0.8, NAVY)
    tb = add_textbox(s, 0.5, 0.1, 12, 0.6)
    p = tb.text_frame.paragraphs[0]
    p.text = "RISK ASSESSMENT MATRIX"; p.font.size = PP(28); p.font.bold = True; p.font.color.rgb = PC(*WHITE)

    risks = [
        ("T-01", "JWT secret default in prod", "Medium", "Critical", "P0", RED),
        ("T-02", "PII stored plaintext", "High", "Critical", "P0", RED),
        ("T-03", "SSRF via resume URLs", "Medium", "High", "P1", ORANGE),
        ("M-01", "SBERT 256-token truncation", "High", "Medium", "P1", ORANGE),
        ("M-03", "Simulated demographics", "Certain", "High", "P1", ORANGE),
        ("B-01", "GDPR non-compliance (PII)", "High", "Critical", "P0", RED),
        ("B-03", "Biased outcomes from data", "Medium", "Critical", "P0", RED),
        ("T-05", "Scoring weight inconsistency", "Certain", "Medium", "P1", ORANGE),
    ]
    y = 1.1
    add_rect(s, 0.5, y, 12.3, 0.5, NAVY)
    for val, cx in [("ID", 0.7), ("Risk", 1.8), ("Likelihood", 6.5), ("Impact", 8.2), ("Priority", 10.0)]:
        tb = add_textbox(s, cx, y + 0.02, 2, 0.45)
        p = tb.text_frame.paragraphs[0]; p.text = val; p.font.size = PP(10); p.font.bold = True
        p.font.color.rgb = PC(*WHITE)
    y += 0.55
    for rid, rdesc, rlike, rimpact, rprio, rcolor in risks:
        bg = (255, 240, 240) if rcolor == RED else (255, 248, 235)
        add_rect(s, 0.5, y, 12.3, 0.5, bg)
        add_rect(s, 0.5, y, 0.1, 0.5, rcolor)
        for val, cx in [(rid, 0.7), (rdesc, 1.8), (rlike, 6.5), (rimpact, 8.2), (rprio, 10.0)]:
            tb = add_textbox(s, cx, y + 0.02, 4.5, 0.45)
            p = tb.text_frame.paragraphs[0]; p.text = val; p.font.size = PP(9)
            if val in ("P0", "Critical"):
                p.font.color.rgb = PC(*RED); p.font.bold = True
            elif val in ("P1", "High"):
                p.font.color.rgb = PC(*ORANGE); p.font.bold = True
            else:
                p.font.color.rgb = PC(*DGRAY)
        y += 0.55

    # ── 11. Roadmap Slide ────────────────────────────
    s = add_slide()
    add_bg(s, WHITE)
    add_rect(s, 0, 0, 13.333, 0.8, DARK_BLUE)
    tb = add_textbox(s, 0.5, 0.1, 12, 0.6)
    p = tb.text_frame.paragraphs[0]
    p.text = "ENTERPRISE ENHANCEMENT ROADMAP"; p.font.size = PP(24); p.font.bold = True; p.font.color.rgb = PC(*WHITE)

    phases = [
        ("Phase 1: Foundation (Q1-Q2)", GREEN, [
            "Skill Graph Knowledge Base (Neo4j / NetworkX)",
            "SHAP Explainability Integration",
            "Bias Transparency Report Generator",
        ]),
        ("Phase 2: Advanced ML (Q3-Q4)", STEEL, [
            "Cross-Encoder Reranking Pipeline",
            "Interview Prediction Model (XGBoost)",
            "LLM-Based Candidate Summary Generator",
        ]),
        ("Phase 3: Enterprise Scale (Q5-Q6)", NAVY, [
            "Multi-Tenant SaaS Architecture",
            "Internal Mobility Engine",
            "Recruiter AI Copilot",
        ]),
    ]
    y = 1.2
    for phase_title, color, items in phases:
        add_rect(s, 1, y, 11.333, 0.55, color)
        tb = add_textbox(s, 1.3, y + 0.02, 10, 0.5)
        p = tb.text_frame.paragraphs[0]; p.text = phase_title
        p.font.size = PP(14); p.font.bold = True; p.font.color.rgb = PC(*WHITE)
        y += 0.6
        for item in items:
            add_rect(s, 1.3, y, 0.08, 0.45, color)
            tb = add_textbox(s, 1.7, y, 10, 0.45)
            p = tb.text_frame.paragraphs[0]; p.text = item; p.font.size = PP(12); p.font.color.rgb = PC(*DGRAY)
            y += 0.5
        y += 0.2

    # ── 12. Immediate Actions Slide ──────────────────
    s = add_slide()
    add_bg(s, WHITE)
    add_rect(s, 0, 0, 13.333, 0.8, RED)
    tb = add_textbox(s, 0.5, 0.1, 12, 0.6)
    p = tb.text_frame.paragraphs[0]
    p.text = "IMMEDIATE ACTIONS (Sprint 1-2)"; p.font.size = PP(24); p.font.bold = True; p.font.color.rgb = PC(*WHITE)

    actions = [
        ("P0", "Fix JWT secret enforcement — block startup with weak secrets", "2h"),
        ("P0", "Integrate EncryptionManager — encrypt PII at rest", "2-3d"),
        ("P0", "Fix assessment submission bug (wrong collection query)", "30min"),
        ("P0", "Remove tlsAllowInvalidCertificates=true", "15min"),
        ("P1", "Fix SSRF in resume_tasks.py — add URL validation", "1d"),
        ("P1", "Unify scoring weights into canonical ScoringConfig", "1d"),
        ("P1", "Activate RBAC decorators in route files", "2d"),
    ]
    y = 1.1
    for prio, action, effort in actions:
        color = RED if prio == "P0" else ORANGE
        shape = add_rect(s, 1, y + 0.05, 0.8, 0.5, color)
        shape.text_frame.paragraphs[0].text = prio
        shape.text_frame.paragraphs[0].font.size = PP(10)
        shape.text_frame.paragraphs[0].font.bold = True
        shape.text_frame.paragraphs[0].font.color.rgb = PC(*WHITE)
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        tb = add_textbox(s, 2.1, y + 0.05, 8.5, 0.5)
        p = tb.text_frame.paragraphs[0]; p.text = action; p.font.size = PP(12); p.font.color.rgb = PC(*DGRAY)
        tb = add_textbox(s, 10.8, y + 0.05, 2, 0.5)
        p = tb.text_frame.paragraphs[0]; p.text = effort; p.font.size = PP(11)
        p.font.color.rgb = PC(*STEEL); p.font.italic = True
        y += 0.62

    # ── 13. Enterprise Readiness Slide ───────────────
    s = add_slide()
    add_bg(s, WHITE)
    add_rect(s, 0, 0, 13.333, 0.8, NAVY)
    tb = add_textbox(s, 0.5, 0.1, 12, 0.6)
    p = tb.text_frame.paragraphs[0]
    p.text = "ENTERPRISE READINESS RATING"; p.font.size = PP(28); p.font.bold = True; p.font.color.rgb = PC(*WHITE)

    ratings = [
        ("AI Innovation", 5, GREEN),
        ("Fairness Engineering", 5, GREEN),
        ("Security Hardening", 2, RED),
        ("Production Scalability", 3, YELLOW),
        ("MLOps Maturity", 2, RED),
        ("Compliance Integrity", 4, GREEN),
    ]
    y = 1.2
    for name, score, color in ratings:
        add_rect(s, 1.5, y, 4.5, 0.6, (245, 248, 252))
        tb = add_textbox(s, 1.7, y + 0.05, 4, 0.5)
        p = tb.text_frame.paragraphs[0]; p.text = name; p.font.size = PP(14); p.font.color.rgb = PC(*DGRAY)
        # Score bar
        bar_max = 5.0
        bar_w = (score / 5) * bar_max
        add_rect(s, 6.5, y + 0.1, bar_w, 0.4, color)
        add_rect(s, 6.5, y + 0.1, bar_max, 0.4, (230, 230, 230))  # background
        add_rect(s, 6.5, y + 0.1, bar_w, 0.4, color)  # filled
        tb = add_textbox(s, 6.5 + bar_max + 0.3, y + 0.05, 1.5, 0.5)
        p = tb.text_frame.paragraphs[0]; p.text = f"{score}/5"; p.font.size = PP(14); p.font.bold = True
        p.font.color.rgb = PC(*color)
        y += 0.75

    # Overall score
    tb = add_textbox(s, 3.5, y + 0.5, 6, 1.0)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "Overall Enterprise Readiness: 7.4 / 10"
    p.font.size = PP(22); p.font.bold = True; p.font.color.rgb = PC(*NAVY)
    p.alignment = PP_ALIGN.CENTER

    # ── 14. Top 3 ROI Slide ──────────────────────────
    s = add_slide()
    add_bg(s, NAVY)
    tb = add_textbox(s, 1, 0.8, 11, 1.0)
    p = tb.text_frame.paragraphs[0]
    p.text = "TOP 3 HIGHEST-ROI INVESTMENTS"; p.font.size = PP(30); p.font.bold = True; p.font.color.rgb = PC(*WHITE)
    p.alignment = PP_ALIGN.CENTER

    rois = [
        ("1", "Security Hardening", "PII encryption + JWT enforcement + RBAC activation\nEliminates compliance risk with moderate effort"),
        ("2", "Scoring Weight Unification", "Resolves the most impactful consistency issue\naffecting candidate outcomes across the platform"),
        ("3", "Cross-Encoder Reranking", "Delivers 15-25% ranking quality improvement\nfor the ML pipeline with minimal infrastructure change"),
    ]
    y = 2.2
    for num, title, desc in rois:
        # Number circle
        shape = add_rect(s, 1.5, y, 0.8, 0.8, STEEL)
        shape.text_frame.paragraphs[0].text = num
        shape.text_frame.paragraphs[0].font.size = PP(24)
        shape.text_frame.paragraphs[0].font.bold = True
        shape.text_frame.paragraphs[0].font.color.rgb = PC(*WHITE)
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        # Title
        tb = add_textbox(s, 2.8, y, 9, 0.45)
        p = tb.text_frame.paragraphs[0]; p.text = title
        p.font.size = PP(18); p.font.bold = True; p.font.color.rgb = PC(*WHITE)
        # Description
        tb = add_textbox(s, 2.8, y + 0.45, 9, 0.6)
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = desc
        p.font.size = PP(12); p.font.color.rgb = PC(180, 200, 220)
        y += 1.5

    # ── 15. Closing Slide ────────────────────────────
    s = add_slide()
    add_bg(s, NAVY)
    add_rect(s, 0, 3.4, 13.333, 0.06, STEEL)
    tb = add_textbox(s, 2, 1.5, 9, 2.0)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "Thank You"
    p.font.size = PP(44); p.font.bold = True; p.font.color.rgb = PC(*WHITE)
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph(); p.text = "Smart Hiring Enterprise Audit Report v1.0"
    p.font.size = PP(16); p.font.color.rgb = PC(*STEEL); p.alignment = PP_ALIGN.CENTER

    tb2 = add_textbox(s, 2, 4.2, 9, 2.0)
    tf2 = tb2.text_frame; tf2.word_wrap = True
    for txt in ["CONFIDENTIAL — INTERNAL AUDIT",
                "Distribution: CTO, VP Engineering, ML Lead, Security Team",
                "Next Review: May 2026",
                "Enterprise Architecture & AI Governance Team"]:
        p = tf2.add_paragraph(); p.text = txt
        p.font.size = PP(12); p.font.color.rgb = PC(150, 170, 190); p.alignment = PP_ALIGN.CENTER

    prs.save(str(out_path))
    print(f"  PPTX saved: {out_path}")
    sz = os.path.getsize(out_path)
    print(f"  Size: {sz / 1024:.0f} KB  |  Slides: {len(prs.slides)}")


# ═══════════════════════════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════════════════════════
if __name__ == "__main__":
    print("=" * 65)
    print("  Smart Hiring Enterprise Audit — Document Generator")
    print("  Generating: PDF  |  DOCX  |  PPTX")
    print("=" * 65)

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    DIAG_DIR.mkdir(parents=True, exist_ok=True)

    # 1. Parse markdown
    print("\n[1/5] Parsing audit report markdown...")
    elements = parse_markdown(MD_FILE)
    print(f"  Parsed {len(elements)} elements, {len(MERMAID_BLOCKS)} Mermaid diagrams found")

    # 2. Render diagrams
    print("\n[2/5] Rendering Mermaid diagrams via kroki.io...")
    diag_paths = render_all_diagrams()
    print(f"  {len(diag_paths)} diagrams processed")

    # 3. Generate PDF
    print("\n[3/5] Generating PDF...")
    generate_pdf(elements, diag_paths, OUTPUT_DIR / "Smart_Hiring_Enterprise_Audit.pdf")

    # 4. Generate DOCX
    print("\n[4/5] Generating DOCX...")
    generate_docx(elements, diag_paths, OUTPUT_DIR / "Smart_Hiring_Enterprise_Audit.docx")

    # 5. Generate PPTX
    print("\n[5/5] Generating PPTX...")
    generate_pptx(elements, diag_paths, OUTPUT_DIR / "Smart_Hiring_Enterprise_Audit.pptx")

    print("\n" + "=" * 65)
    print("  ALL DOCUMENTS GENERATED SUCCESSFULLY")
    print(f"  Output directory: {OUTPUT_DIR}")
    print("=" * 65)
    # List files
    for f in sorted(OUTPUT_DIR.glob("*.*")):
        if f.is_file() and f.suffix in (".pdf", ".docx", ".pptx"):
            print(f"  {f.name}  ({os.path.getsize(f)/1024:.0f} KB)")
